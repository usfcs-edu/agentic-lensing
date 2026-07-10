#!/usr/bin/env python3
"""
Structural lint for the progress deck. Re-opens the saved .pptx and asserts what
can be checked without rendering.

What this CAN catch:
  1. a slide with no speaker notes (the provenance contract),
  2. a shape whose declared geometry runs off the slide,
  3. a picture stretched away from its native pixel aspect,
  4. a picture whose media failed to embed.

What this CANNOT catch, by construction:
  - table row heights are advisory; the renderer grows rows to fit text, so a
    table can overflow downward while its declared geometry is in bounds,
  - python-pptx has no working autofit, so a textbox's text can draw outside the
    box with the box itself still in bounds.
Those need the visual render:
    soffice --headless --convert-to pdf <pptx> && pdftoppm -png -r 80 <pdf> page

Run:
    python3 progress/2026-07-09/lint_deck.py [deck.pptx]
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

TOL = Inches(0.02)
ASPECT_TOL = 0.02          # 2% relative deviation from native aspect


def lint(path: Path) -> int:
    prs = Presentation(str(path))
    W, H = prs.slide_width, prs.slide_height
    problems: list[str] = []
    n_pics = 0

    for i, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide:
            problems.append(f"slide {i}: no notes slide")
        elif not slide.notes_slide.notes_text_frame.text.strip():
            problems.append(f"slide {i}: empty speaker notes")

        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            name = f"slide {i} <{shape.shape_type}>"
            if shape.left < -TOL or shape.top < -TOL:
                problems.append(f"{name}: off top/left "
                                f"({shape.left/914400:.2f}, {shape.top/914400:.2f} in)")
            if shape.width and shape.left + shape.width > W + TOL:
                over = (shape.left + shape.width - W) / 914400
                problems.append(f"{name}: {over:.2f}in past the right edge")
            if shape.height and shape.top + shape.height > H + TOL:
                over = (shape.top + shape.height - H) / 914400
                problems.append(f"{name}: {over:.2f}in past the bottom edge")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_pics += 1
                try:
                    px_w, px_h = Image.open(io.BytesIO(shape.image.blob)).size
                except Exception as exc:                       # noqa: BLE001
                    problems.append(f"{name}: embedded image failed to decode ({exc})")
                    continue
                native = px_w / px_h
                placed = shape.width / shape.height
                if abs(placed - native) / native > ASPECT_TOL:
                    problems.append(
                        f"{name}: aspect distorted, native {native:.3f} "
                        f"vs placed {placed:.3f}")

    print(f"{path.name}: {len(prs.slides)} slides, {n_pics} pictures")
    if problems:
        print(f"\n{len(problems)} problem(s):")
        for p in problems:
            print(f"  - {p}")
        return 1
    print("lint clean: notes on every slide, all shapes in bounds, "
          "no distorted pictures.")
    return 0


if __name__ == "__main__":
    default = Path(__file__).resolve().parent / "claudenet_v4_dr11_2026-07-09.pptx"
    sys.exit(lint(Path(sys.argv[1]) if len(sys.argv) > 1 else default))
