#!/usr/bin/env python3
"""
Build the "ClaudeNet v4 + the DR11-South Sweep" progress deck.

21 slides, 16:9, for the internal research group. Positions ClaudeNet v4 against
the Huang-lineage ResNet/EfficientNet baselines and their sweeps, reports what the
DR11-south sweep found, states the confidence level, and lays out future work.

Layout helpers are copied (not imported) from tools/spectrumfm/build_deck.py --
the same way progress/2026-05-29/build_deck.py copied them -- so this deck stays
hermetic and regenerating it can never perturb another tracked artifact. Three
helpers are added here that the precedent lacks: new_slide() (forces speaker
notes), add_picture_fit() (letterboxes, since the embedded PNGs run from 1.00 to
2.30 aspect and none are 16:9), and fits_body() (a line budget, because
python-pptx has no working autofit and silently draws text outside its box).

PROVENANCE. The v4 result artifacts (resweep_v4_summary.json,
dr11_finetune_gate.json, survivors_dr11s_v4.parquet) live on Perlmutter scratch
and are not in this repo, so every v4 number here is transcribed from
    reproductions/claudenet/papers/v4_section.tex
    reproductions/dr11-campaign-v4/papers/main.tex
Each slide's speaker notes name its source.

Run:
    python3 progress/2026-07-09/build_deck.py     # python3 == ~/.venvs/gdb/bin/python3
Output:
    progress/2026-07-09/claudenet_v4_dr11_2026-07-09.pptx
"""
from __future__ import annotations

import math
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
REPO_ROOT = HERE.parents[1]
REPRO = REPO_ROOT / "reproductions"
FIG = HERE / "figures"
CN_FIG = REPRO / "claudenet" / "papers" / "figures"
D11_FIG = REPRO / "dr11-campaign" / "papers" / "figures"
D11V4_FIG = REPRO / "dr11-campaign-v4" / "papers" / "figures"
OUTPUT = HERE / "claudenet_v4_dr11_2026-07-09.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
LEFT_MARGIN = Inches(0.5)
TOP_MARGIN = Inches(0.35)
CONTENT_W = SLIDE_W - 2 * LEFT_MARGIN
CONTENT_W_IN = 12.33

NAVY = RGBColor(0x18, 0x2C, 0x5B)
ACCENT = RGBColor(0xC6, 0x3A, 0x3A)
GOOD = RGBColor(0x1E, 0x7A, 0x3C)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)


# ---------------------------------------------------------------- helpers ----
def set_text(textframe, text, size=24, bold=False, color=INK, align=None):
    textframe.clear()
    p = textframe.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def new_slide(prs, notes: str):
    """Blank slide with speaker notes set at birth.

    slide.notes_slide lazily *creates* the notes part on access, so routing every
    slide through this factory is what makes "notes on every slide" structurally
    true rather than a per-slide chore that silently regresses.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_title(slide, text, size=30):
    box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, CONTENT_W, Inches(0.75))
    set_text(box.text_frame, text, size=size, bold=True, color=NAVY)
    return box


def add_subtitle(slide, text, top=Inches(1.08), size=15, color=MUTED, width=None,
                 left=LEFT_MARGIN):
    box = slide.shapes.add_textbox(left, top, width or CONTENT_W, Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    set_text(tf, text, size=size, color=color)
    return box


def _bullet_prefix(level: int, text: str) -> str:
    if not text:
        return ""                              # spacer paragraph, no orphan bullet
    if level > 0:
        return "– "
    if len(text) > 1 and text[0].isdigit() and text[1] == ".":
        return ""                              # already numbered; don't double-mark
    return "• "


def add_body(slide, top, height, bullets, size=17, indent_size=14,
             left=LEFT_MARGIN, width=CONTENT_W):
    """bullets: list[str] or list[tuple(level, str)]. "" gives a spacer line."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE      # be explicit; pptx cannot autofit anyway
    first = True
    for item in bullets:
        level, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = _bullet_prefix(level, text) + text
        run.font.size = Pt(size if level == 0 else indent_size)
        run.font.color.rgb = INK
    return box


def add_takeaway(slide, text, top=Inches(6.62), color=ACCENT, size=14):
    box = slide.shapes.add_textbox(LEFT_MARGIN, top, CONTENT_W, Inches(0.52))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def add_footer(slide, text):
    box = slide.shapes.add_textbox(LEFT_MARGIN, SLIDE_H - Inches(0.42),
                                   CONTENT_W, Inches(0.3))
    set_text(box.text_frame, text, size=10, color=MUTED)


def add_table(slide, top, height, rows, col_widths_in, left=LEFT_MARGIN,
              header_row=True, size=12):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top,
                                   Inches(sum(col_widths_in)), height)
    table = shape.table
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            tf = table.cell(r, c).text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(size)
            if r == 0 and header_row:
                run.font.bold = True
                run.font.color.rgb = NAVY
            else:
                run.font.color.rgb = INK
    return shape


def add_picture_fit(slide, img_path, box_left, box_top, box_w, box_h):
    """Letterbox an image into a box: preserve aspect, center, never distort.

    add_picture(p, l, t, width=W) scales height proportionally, but pinning both
    width and height to a box that isn't the image's aspect stretches it. The
    embedded PNGs span aspect 1.00-2.30, so the fit has to be computed.
    """
    img_path = Path(img_path)
    if not img_path.exists():                       # fail loudly at build time
        raise FileNotFoundError(img_path)
    px_w, px_h = Image.open(img_path).size
    ar = px_w / px_h
    if box_w / box_h > ar:                          # box is wider -> height-bound
        h, w = box_h, int(round(box_h * ar))
    else:                                           # width-bound
        w, h = box_w, int(round(box_w / ar))
    left = int(box_left + (box_w - w) // 2)
    top = int(box_top + (box_h - h) // 2)
    return slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)


def est_lines(text: str, box_w_in: float, font_pt: float) -> int:
    """Wrapped-line estimate: average glyph advance ~= 0.5 em."""
    chars_per_line = max(1, int(box_w_in * 72.0 / (font_pt * 0.5)))
    return max(1, math.ceil(len(text) / chars_per_line))


def fits_body(bullets, box_w_in: float, box_h_in: float, font_pt: float,
              indent_pt: float = 14.0, line_factor: float = 1.28) -> bool:
    total = 0.0
    for item in bullets:
        level, text = item if isinstance(item, tuple) else (0, item)
        pt = font_pt if level == 0 else indent_pt
        total += est_lines(text, box_w_in, pt) * pt * line_factor / 72.0
    return total <= box_h_in


def body(slide, top_in, height_in, bullets, size=17, left_in=0.5, width_in=CONTENT_W_IN):
    """add_body + an enforced line budget, so an over-stuffed slide fails the build."""
    if not fits_body(bullets, width_in, height_in, size):
        raise ValueError(f"body overflows {width_in}x{height_in}in @ {size}pt: "
                         f"{bullets[0]!r}...")
    return add_body(slide, Inches(top_in), Inches(height_in), bullets, size=size,
                    left=Inches(left_in), width=Inches(width_in))


CITE = ("Provenance: v4 numbers are transcribed from "
        "reproductions/claudenet/papers/v4_section.tex and "
        "reproductions/dr11-campaign-v4/papers/main.tex. The v4 result artifacts "
        "(resweep_v4_summary.json, dr11_finetune_gate.json, "
        "survivors_dr11s_v4.parquet) live on Perlmutter scratch, not in this repo.")

NEW3 = ["327526_9059", "340971_2860", "325163_2981"]


# ----------------------------------------------------------------- slides ----
def slide_01_cover(prs):
    slide = new_slide(prs, notes=(
        "Progress update for the group, 9 July 2026.\n\n"
        "Cutouts are the three genuinely-new systems from the v3blend8 DR11-south "
        "campaign after LensJudge + HSC PDR3 tier-2 vetting: s_327526_9059, "
        "s_340971_2860, s_325163_2981. Top row DESI grz (0.262\"/px), bottom row "
        "HSC grizy (0.168\"/px). Source: reproductions/dr11-campaign/papers/main.tex, "
        "Table tab:confirmed and Appendix app:gallery.\n\n"
        "Frame the talk with the one-liner: the recall gains are large and "
        "confidence-interval'd; the discovery yield is small and resolution-bounded."))
    box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.75), CONTENT_W, Inches(1.5))
    set_text(box.text_frame, "ClaudeNet v4 and the DR11-South Sweep",
             size=38, bold=True, color=NAVY)
    b2 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.72), CONTENT_W, Inches(0.9))
    tf = b2.text_frame
    tf.word_wrap = True
    set_text(tf, "What we improved, what the sweep discovered, and how confident we are",
             size=19, color=MUTED)
    b3 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.42), CONTENT_W, Inches(0.4))
    set_text(b3.text_frame, "Greg Benson, USF CS  ·  9 July 2026  ·  progress update",
             size=14, color=INK)

    cap = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.05), CONTENT_W, Inches(0.35))
    set_text(cap.text_frame,
             "The three genuinely new lens systems from the DR11-south campaign "
             "— DESI grz (top) and HSC grizy (bottom)",
             size=12, color=MUTED)

    for i, sysid in enumerate(NEW3):
        center = 3.0 + i * 3.665                       # three evenly spaced columns
        x = Inches(center - 0.75)
        add_picture_fit(slide, D11_FIG / f"cand_s_{sysid}_desi.png",
                        x, Inches(3.42), Inches(1.5), Inches(1.5))
        add_picture_fit(slide, D11_FIG / f"cand_s_{sysid}_hsc.png",
                        x, Inches(4.97), Inches(1.5), Inches(1.5))
        lab = slide.shapes.add_textbox(Inches(center - 1.1), Inches(6.52),
                                       Inches(2.2), Inches(0.32))
        set_text(lab.text_frame, f"s_{sysid}", size=12, color=INK,
                 align=PP_ALIGN.CENTER)

    add_footer(slide, "53,809,040 galaxies swept  ·  3 confirmed new systems  ·  "
                      "134,078 v4 candidates pending vetting")


def slide_02_bluf(prs):
    slide = new_slide(prs, notes=(
        "Bottom line up front. Hold the two halves apart all talk: recall is "
        "measured and confidence-interval'd; discovery is small and bounded by "
        "pixel resolution, not by the model.\n\n"
        "Numbers: 0.54 -> 0.87 Inchausti grade-A is dr11-campaign-v4/papers/main.tex "
        "Table 2. The 3 new systems are dr11-campaign/papers/main.tex (post-hoc "
        "literature crosscheck of 8 distinct new-class systems left 3 genuinely new)."
        f"\n\n{CITE}"))
    add_title(slide, "Bottom line")
    body(slide, 1.25, 5.05, [
        "The Huang lineage's ceiling was never architecture: a 194k-param shielded "
        "ResNet matches a 20.5M-param EfficientNetV2-S within ±0.003 AUC, and the "
        "published meta-learner collapsed to a simple average.",
        "ClaudeNet's four real levers, in the order each exposed the next: "
        "diversity (v1) → honest deployment-scale negatives (v2) → contaminant "
        "awareness (v3) → release-native adaptation (v4).",
        "v4 = v3's contaminant-aware members + two additive, release-portable levers: "
        "a calibrated-mean stage-1 selector (free) and a DR11-native warm-start "
        "fine-tune (one short GPU pass).",
        "It is the best-recall configuration ever measured on DR11-south: "
        "Inchausti grade-A 0.54 → 0.87, Storfer grade-A (hard) 0.32 → 0.825.",
        "But: the v3 DR10 sweep produced zero net-new lenses; the v3 DR11 campaign "
        "produced three after HSC vetting; and v4's 134,078 new candidates are not "
        "yet vetted at all.",
        "On the one independent, unbiased test — Euclid grade-A-vs-C AUC — all four "
        "model generations tie within one standard error.",
    ], size=17)
    add_takeaway(slide, "Recall went up a lot, and we can prove it. Discovery went up "
                        "a little, and only where higher-resolution pixels exist.")


def slide_03_lineage(prs):
    slide = new_slide(prs, notes=(
        "The baselines we are positioning against.\n\n"
        "Huang+2020: Lanusse-2018 CMU DeepLens ResNet-46, 3,508,833 params; DECaLS DR7, "
        "6,242,507 galaxies; 342 candidates (60A/106B/176C).\n"
        "Huang+2021: 'shielded' ResNet, 59,905 params (58.6x reduction), deployed as a "
        "2-model ensemble; DR8 ~14,000 deg^2, 17,290,814 parent; 1,312 candidates "
        "(216A/199B/897C), 1,210 new.\n"
        "Storfer+2024: same shielded ResNet at 1,961 training lenses; DR9 ~19,000 deg^2, "
        "45.26M cutouts; 1,895 candidates, 1,512 new.\n"
        "Inchausti+2025: shielded ResNet (194,501) + EfficientNetV2-S (20,543,145) + "
        "FWLS meta (1,201); DR10 ~14,000 deg^2, ~43M cutouts; 811 new candidates.\n\n"
        "All four vetted by by-eye visual grading. The chart's second panel is "
        "deliberately not commensurable across regimes -- say so out loud.\n"
        "Sources: reproductions/{huang-2020,huang-2021,inchausti-2025}/README.md."))
    add_title(slide, "Where we started: the Huang lineage and its sweeps")
    add_picture_fit(slide, FIG / "fig1_lineage_sweeps.png",
                    Inches(1.05), Inches(1.15), Inches(11.2), Inches(5.35))
    add_takeaway(slide, "Eight years, 8.6× more galaxies scanned, and candidate yield "
                        "did not follow — while the vetting bar rose from by-eye grading "
                        "to HSC tier-2 confirmation.")


def slide_04_negative_result(prs):
    slide = new_slide(prs, notes=(
        "The negative result that started ClaudeNet. Reproduced controlled AUC from "
        "reproductions/inchausti-2025/ (their Fig 6): ResNet 0.9984, EffNet 0.9987, "
        "meta 0.9989 -- and the simple average is also 0.9989.\n\n"
        "The meta-learner is a 1,201-param FWLS MLP: Linear(2->300)->ReLU->Linear(300->1) "
        "on [p_shielded, p_effnet]. It collapsed to averaging because both base models "
        "were trained on byte-identical data, so their errors correlate at ~1.0 and a "
        "combiner has nothing to exploit.\n\n"
        "Stage-D honest baseline (matched-FPR recovery): meta @1% FPR = 0.908 Storfer / "
        "0.968 Inchausti; @0.1% = 0.755 / 0.845. Inchausti's own finding: the "
        "negative:positive ratio, not architecture or AUC, sets usability at a real "
        "operating point."))
    add_title(slide, "The negative result that started ClaudeNet")
    add_subtitle(slide, "Reproduced, controlled comparison — architecture is not the bottleneck")
    add_table(slide, Inches(1.62), Inches(2.4), [
        ["model", "params (reproduced)", "test AUC", "role"],
        ["shielded ResNet", "194,501", "0.9984", "base model 1"],
        ["EfficientNetV2-S", "20,543,145", "0.9987", "base model 2"],
        ["FWLS meta-learner", "1,201", "0.9989", "stacked combiner"],
        ["simple average", "—", "0.9989", "the meta-learner's actual behaviour"],
    ], col_widths_in=[3.0, 3.0, 2.2, 4.1], size=13)
    body(slide, 4.35, 2.0, [
        "A 194k-param ResNet ties a 20.5M-param EfficientNet to within ±0.003 AUC — "
        "105× the parameters, no gain.",
        "The combiner had nothing to exploit: both bases saw byte-identical data, so "
        "their errors are correlated at ~1.0.",
        "Honest Stage-D baseline to beat (recovery @ matched FPR): "
        "meta @1% = 0.908 Storfer / 0.968 Inchausti;  @0.1% = 0.755 / 0.845.",
    ], size=16)
    add_takeaway(slide, "Bigger backbones and a learned combiner bought nothing, because "
                        "every model saw the same data and made the same mistakes. "
                        "The levers are elsewhere: diversity, negatives, contaminants, domain.")


def slide_05_metric(prs):
    slide = new_slide(prs, notes=(
        "The metric. Recovery @ matched false-positive rate, using arithmetic copied "
        "verbatim from reproductions/inchausti-2025/22_fpr_operating_point.py into "
        "claudenet/_ensemble.py, so every number in this deck is directly comparable "
        "to the reproduced baseline.\n\n"
        "Why not AUC: AUC saturates near 1.0 (all four models above sit at 0.998-0.999) "
        "and averages over the whole ROC. A real sweep operates at a fixed, tiny FPR "
        "where only the tail matters. A 45M-cutout sweep at 1% FPR yields ~450k false "
        "positives; scans run near 0.01%.\n\n"
        "Figure: claudenet/papers/figures/flagship_operating_point.png."))
    add_title(slide, "The metric that makes every number here comparable")
    add_subtitle(slide, "Recovery @ matched false-positive rate — the exact arithmetic of "
                        "inchausti-2025/22_fpr_operating_point.py, reused verbatim")
    body(slide, 1.72, 3.1, [
        "AUC saturates near 1.0 and hides the tail; a sweep lives entirely in the tail.",
        "A 45M-cutout sweep at 1% FPR means ~450,000 false positives. Real scans run "
        "near 0.01% FPR.",
        "So we report recovery at 1%, 0.1% and 0.01% FPR against held-out Storfer and "
        "Inchausti graded lenses.",
        "Same held-out catalogs throughout — and from v4 on, they are excluded from "
        "training as both positives and negatives.",
    ], size=16, width_in=6.1)
    add_picture_fit(slide, CN_FIG / "flagship_operating_point.png",
                    Inches(6.95), Inches(1.6), Inches(5.85), Inches(4.75))
    add_takeaway(slide, "One number, one operating point: recovery at matched FPR. "
                        "Everything that follows is apples-to-apples against the "
                        "reproduced baseline.")


def slide_06_v1_diversity(prs):
    slide = new_slide(prs, notes=(
        "ClaudeNet v1. Five to six deliberately decorrelated members: EfficientNetV2-S "
        "x2, EfficientNet-B3, an AION-1 frozen-embedding MLP probe, the shielded-194k "
        "ResNet, and the Lanusse ResNet-46. Each isotonic-calibrated, then combined.\n\n"
        "Member score correlation fell from the lineage's ~1.0 to Pearson ~0.31 / "
        "Spearman ~0.45. Beat the published meta on 4/4 matched-FPR cells.\n\n"
        "Other v1 phases: conformal selection certified FDR control; deep-ensemble "
        "triage cut selective error 0.022 -> 0.0002 at 50% coverage; test-time D4 "
        "pooling added +0.041 mean. Domain adaptation was a negative result: the "
        "north/south gap is real (north@1% 0.688 vs south 0.874) but naive MMD at "
        "lambda=1 hurt both domains by ~0.06.\n\n"
        "Caveat to state: v1's 0.1% threshold was pinned by only ~6-7 of ~6,500 held-out "
        "negatives -- about +/-10pp of noise, no CI. v2 fixes exactly this.\n"
        "Figure: claudenet/papers/figures/diversity_heatmap.png."))
    add_title(slide, "v1 — diversity, not combiner cleverness, is the lever")
    body(slide, 1.25, 5.0, [
        "Five decorrelated members (EfficientNetV2-S ×2, EfficientNet-B3, an AION-1 "
        "frozen-embedding probe, shielded-194k ResNet, Lanusse ResNet-46).",
        "Member score correlation: Pearson ~0.31 / Spearman ~0.45, against the "
        "lineage's collapsed ~1.0.",
        "Beat the published meta-learner on 4/4 matched-FPR cells — with a naive "
        "average, no clever combiner:",
        (1, "Storfer @1% 0.908→0.938,  @0.1% 0.755→0.853"),
        (1, "Inchausti @1% 0.968→0.980,  @0.1% 0.845→0.935"),
        "Also shipped: certified-FDR conformal selection, and deep-ensemble triage "
        "(selective error 0.022 → 0.0002 at 50% coverage).",
        "Honest caveat: the 0.1% threshold was pinned by ~6–7 of only ~6,500 held-out "
        "negatives — roughly ±10 pp of noise, no confidence interval.",
    ], size=15, width_in=6.4)
    add_picture_fit(slide, CN_FIG / "diversity_heatmap.png",
                    Inches(7.35), Inches(1.35), Inches(5.4), Inches(5.0))
    add_takeaway(slide, "With decorrelated members even a plain average beats the "
                        "published correlated-base meta-learner — and the gap is widest "
                        "at the strict thresholds that control purity.")


def slide_07_v2(prs):
    slide = new_slide(prs, notes=(
        "ClaudeNet v2, the Perlmutter scale-up. NegEval-1M replaced the ~6,500 held-out "
        "negatives with 1,000,000, so the 0.01% FPR point could be measured for the "
        "first time. Deployment-scale hard-negative mining over a 1M pool shipped in "
        "round 1 (+0.157 Inchausti @0.1%); round 2 over-mined and was rejected.\n\n"
        "Two things FAILED their gate and were dropped -- worth saying out loud:\n"
        " - the AION native-griz upgrade (0.535 < 0.647): the synthetic-i 'degradation' "
        "WAS the diversity mechanism (Pearson vs CNN 0.13 degraded vs ~0.65 native).\n"
        " - distillation to a single student: 5.3x faster but -0.111 recovery.\n\n"
        "v2-lean roster: effnet_B, effnet_B3_hard, effnet_S2_hard, resnet46_C_hard, "
        "zoobot_N. Table tab:v2lean, paired bootstrap 10,000 reps.\n\n"
        "Note the v1 collapse at 0.01% Storfer (0.394, below the meta's 0.513): the "
        "degraded-AION member's negative tail poisons the average. The RF combiner was "
        "robust there. This is why v2 dropped AION from the roster."))
    add_title(slide, "v2 — an honest eval, and deployment-scale negatives")
    add_subtitle(slide, "NegEval-1M: 1,000,000 held-out negatives, so the 0.01% FPR "
                        "operating point can be measured at all")
    add_picture_fit(slide, FIG / "fig2_recovery_matched_fpr.png",
                    Inches(1.6), Inches(1.62), Inches(10.1), Inches(4.85))
    add_takeaway(slide, "The gain grows monotonically as the threshold tightens — "
                        "+0.34 Storfer / +0.26 Inchausti at 0.01% FPR — which is exactly "
                        "the regime a 10⁷-galaxy sweep runs in.")


def slide_08_v2_reality(prs):
    slide = new_slide(prs, notes=(
        "The pivot. v2's DR9 sweep: 17,290,814 parent galaxies -> 29,892 survivors -> "
        "group-conformal at FDR<=0.05 selects 1,449 (813 new, 737 new-and-unseen). "
        "Known-lens recall 47.5% in-coverage.\n\n"
        "Then the campaign result: of 601 genuinely-new group-conformal candidates, ZERO "
        "were graded A/B by either independent vetter. The only five graded >=B were "
        "already-catalogued lenses. Rejects were overwhelmingly luminous red galaxies "
        "with companions, rings, or blends.\n\n"
        "That is the diagnosis that produced v3 and v4: the binding constraint is "
        "lens-vs-mimic separation and vetting resolution, not architecture and not AUC. "
        "Source: claudenet/papers/v3_section.tex (motivation), README.md phase 160."))
    add_title(slide, "v2's reality check — the score moved, the discoveries did not")
    body(slide, 1.3, 2.5, [
        "Full DR9 sweep: 17,290,814 parent galaxies → 29,892 survivors → group-conformal "
        "at FDR ≤ 0.05 selects 1,449 (813 new, 737 new-and-unseen).",
        "Known-lens recall into survivors: 47.5% in-coverage. The pipeline works.",
        "Then vetting: of 601 genuinely-new conformal-selected candidates, "
        "zero were graded A/B by either independent vetter.",
        "The rejects were overwhelmingly luminous red galaxies with companions, rings, "
        "or blends — objects that look like lenses at 0.262″/px.",
    ], size=16)

    box = slide.shapes.add_textbox(Inches(1.1), Inches(4.15), Inches(11.1), Inches(1.35))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("Diagnosis: the binding constraint is lens-vs-mimic separation and "
              "vetting resolution — not architecture, and not AUC.")
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "This is the pivot that produced v3 (contaminant-aware training) and v4 " \
              "(release-native adaptation)."
    r2.font.size = Pt(15)
    r2.font.color.rgb = MUTED

    add_takeaway(slide, "A better score produced not one new confirmed lens. "
                        "The wall is telling real lenses from convincing look-alikes — "
                        "and, ultimately, pixel resolution.")


def slide_09_v3(prs):
    slide = new_slide(prs, notes=(
        "ClaudeNet v3, contaminant-aware. New metric: recovery @ matched MIMIC-FPR. "
        "v2-lean, which recovers 96-100% of held-out lenses at random-FPR, recovers only "
        "0.168 Storfer / 0.307 Inchausti at mimic-FPR 0.05 -- the random-FPR metric was "
        "flattering.\n\n"
        "Mimic bank: 148,034 CNN-high DR10 non-lens survivors, morphology-typed via "
        "Tractor, 98.3% non-lens purity by a 120-row agentic gate; 118,901 training rows "
        "+ a frozen 29,724-row held-out mimic-eval.\n\n"
        "Recipe: hold the v2 'hard' recipe byte-identical and swap only the negative "
        "CONTENT -- a fraction f of displaced negatives become typed mimics. f=0.5 "
        "('b50') wins; f=0.7 over-specialises.\n\n"
        "Ship roster v3blend8 = 8 members, keeping BOTH the v2-hard and v3-b50 versions "
        "of each swappable member: effnet_B, zoobot_N, effnet_S2_{hard,b50}, "
        "effnet_B3_{hard,b50}, resnet46_C_{hard,b50}. Source: v3_section.tex tab_a2, "
        "a6_ensemble_refit.json."))
    add_title(slide, "v3 — contaminant-aware training on a typed mimic bank")
    add_subtitle(slide, "Score the models against lens mimics, not random galaxies, "
                        "and the flattering numbers disappear")
    add_table(slide, Inches(1.65), Inches(1.9), [
        ["held-out DR10 mimic-eval", "v2-lean", "v3blend8", ""],
        ["mimic-FPR 0.05 (Storfer)", "0.713", "0.793", "the contaminant metric"],
        ["mimic-FPR 0.01 (Storfer)", "0.460", "0.699", "+0.24 at the strict point"],
        ["random-FPR 0.01 (Storfer)", "0.967", "0.935", "the deliberate cost"],
        ["random-FPR 0.01 (Inchausti)", "0.998", "0.979", "the deliberate cost"],
    ], col_widths_in=[4.0, 2.1, 2.1, 4.1], size=13)
    body(slide, 3.85, 2.5, [
        "Mimic bank: 148,034 CNN-high DR10 non-lens survivors, morphology-typed via "
        "Tractor, 98.3% non-lens purity (118,901 train / 29,724 frozen held-out).",
        "The v2 recipe is held byte-identical; only the negative content changes — "
        "half the displaced negatives become typed mimics (\"b50\").",
        "Ship roster v3blend8 keeps both the v2-hard and the v3-b50 version of each "
        "swappable member — 8 members, near-Pareto on both metrics.",
        "Broad-mimic recovery (DR10-eval): 0.598 original CNN → 0.790 v3blend8 → "
        "0.885 with the v3 head re-ranker (deployed only as a φ=0.05 stage-2 re-rank).",
    ], size=15)
    add_takeaway(slide, "Training against typed mimics roughly halves the hard-contaminant "
                        "miss rate, at a small and deliberate cost on the easy random-FPR metric.")


def slide_10_v3_campaign(prs):
    slide = new_slide(prs, notes=(
        "The v3blend8 DR11-south as-run campaign -- COMPLETE and HSC-vetted. This is the "
        "confirmed catalogue of record.\n\n"
        "53,809,040 DR11-south galaxies swept; per-member 1e-4 union -> 95,104 survivors. "
        "Conformal certified ZERO at this scale (power-limited, not a null result).\n"
        "Top-300 = 102 known + 198 new; discovery set = the top-500 new.\n"
        "LensJudge v3 cascade + HSC PDR3 tier-2 (0.168\"/px, ~1,300 deg^2 ~= 20x Euclid "
        "Q1). Total cost $40.38.\n"
        "Cascade: 500 graded -> 350 escalated -> 26 reached HSC tier-2; tier-1 grades "
        "A/B/C/D = 27/76/125/272; HSC tier-2 grade-A/B = 24 (19 A, 5 B).\n"
        "SuGOHI: 79 SuGOHI HSC committee lenses recovered into survivors at 2.7x the "
        "survivor-median score; 15 of the 24 tier-2 A/B are SuGOHI lenses.\n"
        "Net-new accounting: of the 24, zero are in DESI finder catalogs, 15 are SuGOHI, "
        "9 (8 distinct systems) are new-class. A post-hoc literature crosscheck found 5 "
        "of the 8 previously published -> 3 GENUINELY NEW.\n"
        "Also: 108 DESI-resolution-only A/B (17A/91B), of which 39 genuinely new.\n\n"
        "Source: reproductions/dr11-campaign/papers/main.tex, tab:cascade, tab:confirmed."))
    add_title(slide, "The v3 DR11-south campaign — swept, vetted, and counted honestly")
    add_subtitle(slide, "53,809,040 galaxies → 95,104 union survivors → top-500 new → "
                        "LensJudge v3 cascade → HSC PDR3 tier-2 at 0.168″/px  ($40.38)")
    add_picture_fit(slide, D11_FIG / "desi_hsc_flip.png",
                    Inches(0.55), Inches(1.68), Inches(4.4), Inches(4.4))
    add_picture_fit(slide, D11_FIG / "sugohi_enrichment.png",
                    Inches(5.05), Inches(1.68), Inches(4.4), Inches(4.4))
    body(slide, 1.8, 3.4, [
        "500 graded → 350 escalated → 26 reached HSC tier-2",
        "HSC tier-2 grade-A/B: 24  (19 A, 5 B)",
        "of those 24: 0 in DESI finder catalogs, 15 are SuGOHI lenses",
        "9 new-class detections = 8 distinct systems",
        "post-hoc literature crosscheck: 5 of 8 already published",
        "⇒ 3 genuinely new systems",
        "(+ 39 genuinely new among 108 DESI-resolution-only A/B)",
    ], size=13, left_in=9.5, width_in=3.4)
    add_takeaway(slide, "A full 53.8M-galaxy sweep plus HSC vetting yielded exactly three "
                        "genuinely new lenses. The model recovers known systems well; "
                        "net-new discovery is bounded by resolution.")
    add_footer(slide, "Left: the DESI→HSC resolution flip on the 26 HSC-covered candidates. "
                      "Right: v3blend8 enriches independent SuGOHI lenses to 2.7× the "
                      "survivor-median score.")


def slide_11_v4_overview(prs):
    slide = new_slide(prs, notes=(
        "ClaudeNet v4. Trigger: deploying v3blend8 on DR11-south showed known-lens "
        "grade-A recall at the per-member 1e-4 operating point falling 62% -> 41% "
        "relative to DR10.\n\n"
        "v4 is ADDITIVE. It keeps v3's contaminant-aware members and stacks two "
        "release-portable levers:\n"
        " (1) a combiner change -- free, no retraining;\n"
        " (2) a release-native warm-start fine-tune -- one short GPU pass.\n\n"
        "Both levers are exactly the recipe that transfers to the higher-resolution "
        "surveys that are the program's endpoint: re-harvest the native-resolution "
        "positive pool, mine native hard negatives, warm-start fine-tune, select by mean."
        f"\n\n{CITE}"))
    add_title(slide, "ClaudeNet v4 — two additive, release-portable levers")
    add_subtitle(slide, "Trigger: on DR11-south, v3blend8's grade-A recall at the "
                        "per-member 10⁻⁴ operating point appeared to fall 62% → 41% vs DR10")
    add_picture_fit(slide, FIG / "fig6_v4_pipeline.png",
                    Inches(0.5), Inches(1.6), Inches(8.5), Inches(4.85))
    body(slide, 1.85, 4.4, [
        "v4 is additive — it keeps v3's contaminant-aware members.",
        "",
        "Lever 1 — calibrated-mean stage-1 selector.",
        (1, "Free. No retraining. Removes DR9→DR11 threshold drift entirely."),
        "",
        "Lever 2 — release-native warm-start fine-tune.",
        (1, "One short GPU pass on a 13×-expanded DR11-native positive pool."),
        "",
        "Both levers are the recipe that transfers to the higher-resolution surveys "
        "that are the program's endpoint.",
    ], size=13, left_in=9.25, width_in=3.6)
    add_takeaway(slide, "v4 doesn't replace v3 — it fixes the selection rule and adapts "
                        "the weights to the actual data release.")


def slide_12_lever1(prs):
    slide = new_slide(prs, notes=(
        "Lever 1. The apparent recall collapse is an OPERATING-POINT artifact, not a "
        "model failure.\n\n"
        "Mechanism: the stage-1 selector was a per-member 1e-4 UNION. On the deeper DR11 "
        "imaging the per-member thresholds tighten and resnet46_C saturates (its 1e-4 "
        "threshold moved 0.81 -> 1.00), so the union is dominated by single-member spikes.\n\n"
        "Three facts show the model is intact:\n"
        " (i) 21/90 grade-A lie outside the searched footprint or are cut by the parent "
        "selection, so searchable recall is already 54%, not 41%;\n"
        " (ii) threshold-free AUC of in-parent held-out grade-A vs 2M random DR11 "
        "negatives = 0.9955 for the five-member calibrated mean;\n"
        " (iii) swapping union for the calibrated mean at the SAME survivor budget "
        "recovers grade-A recall 54% -> 75%.\n\n"
        "DENOMINATOR DISCIPLINE -- say this if asked. Those percentages use different "
        "denominators. The authoritative three-configuration comparison is Table 2 of "
        "dr11-campaign-v4/papers/main.tex (position-crossmatch, training-excluded), which "
        "is what slide 14 plots. The plotted curve here uses its own 'held-out in-parent "
        "grade-A' denominator and reads ~0.69 at 95k. Do not mix them on one axis.\n\n"
        "Figure: dr11-campaign-v4/papers/figures/recall_recoverability.png."))
    add_title(slide, "v4 Lever 1 — the recall collapse was the combiner, not the model")
    add_picture_fit(slide, D11V4_FIG / "recall_recoverability.png",
                    Inches(0.5), Inches(1.25), Inches(5.9), Inches(5.15))
    body(slide, 1.25, 2.35, [
        "The union over per-member 10⁻⁴ thresholds degrades on deeper DR11 imaging: "
        "thresholds tighten, resnet46_C saturates (0.81 → 1.00), and the union ends up "
        "dominated by single-member spikes.",
        "Threshold-free AUC of in-parent held-out grade-A vs 2M random DR11 negatives: "
        "0.9955 for the five-member calibrated mean. The model ranks DR11 lenses fine.",
        "The calibrated mean carries no per-release threshold recalibration — it removes "
        "the DR9→DR11 drift entirely, and is now the default stage-1 selector.",
    ], size=13, left_in=6.65, width_in=6.15)

    add_subtitle(slide, "Denominator discipline — three readings, three denominators:",
                 top=Inches(3.78), size=13, color=NAVY,
                 width=Inches(6.15), left=Inches(6.65))
    add_table(slide, Inches(4.3), Inches(2.05), [
        ["reading", "value", "denominator"],
        ["union @10⁻⁴, all grade-A", "41%  (37/90)", "all 90 grade-A"],
        ["union @10⁻⁴, searchable", "54%  (37/69)", "in-parent, in-footprint"],
        ["calibrated mean @95k", "75%", "in-parent"],
        ["calibrated mean @150k", "0.80 (Inch.-A)", "position-crossmatch, held out"],
    ], col_widths_in=[2.35, 1.65, 2.15], left=Inches(6.65), size=11)
    add_takeaway(slide, "Swapping the union for a calibrated mean recovers the recall for "
                        "free — no retraining. The 62%→41% \"collapse\" was a selection "
                        "rule, not a broken model.")


def slide_13_lever2(prs):
    slide = new_slide(prs, notes=(
        "Lever 2. Release-native warm-start fine-tune.\n\n"
        "The v1-v3 positive set was capped at 1,961 DR9-pixel lenses. We re-harvested the "
        "confirmed-lens literature for the DR11-south footprint: VizieR/DES/KiDS/HSC "
        "catalogs, SLACS/BELLS spectroscopic lenses, lensed-quasar compilations, a bulk "
        "SIMBAD object-type pull, and the Euclid-Q1 grade-A/B set. Deduplicated at 5 "
        "arcsec, tiered by confidence -> 3,171 net-new high-confidence (gold+silver) "
        "positives, a ~13x expansion.\n\n"
        "Cutouts re-extracted DR11-south-native (griz, 101px), tier-subsampled to 5,806, "
        "paired with 30k random + 20k hard DR11 negatives. The hard negatives are "
        "CNN-high mimics from the mean-ranked survivor pool, PU-guarded at 10 arcsec and "
        "candidate-region-skimmed (skip the top 2,000, take the next 20,000) so real "
        "unconfirmed lenses do not poison the negative class.\n\n"
        "The three swappable members are warm-started from their _b50 checkpoints and "
        "fine-tuned 12 epochs at lr=3e-4. effnet_B and zoobot_N stay FROZEN as the "
        "random-FPR / diversity anchors.\n\n"
        "Storfer and Inchausti are held out throughout, as positives AND as negatives, so "
        "the gate measures generalisation and not memorisation.\n\n"
        "Gate table = v4_section.tex Table tab:v4gate. AUC does not regress: 0.996->0.998."
        f"\n\n{CITE}"))
    add_title(slide, "v4 Lever 2 — a release-native warm-start fine-tune")
    body(slide, 1.2, 2.6, [
        "v1–v3 trained on 1,961 DR9-pixel lenses. We re-harvested the DR11-south "
        "confirmed-lens literature (VizieR/DES/KiDS/HSC, SLACS/BELLS, lensed quasars, "
        "SIMBAD, Euclid-Q1 A/B), deduped at 5″ and tiered by confidence: "
        "3,171 net-new high-confidence positives — a ~13× expansion.",
        "Cutouts re-extracted DR11-native (griz, 101 px), tier-subsampled to 5,806; "
        "paired with 30k random + 20k hard DR11 negatives (PU-guarded at 10″, and "
        "candidate-region-skimmed so real unconfirmed lenses can't poison the negatives).",
        "Three swappable members warm-started from _b50 and fine-tuned 12 epochs @ "
        "lr 3×10⁻⁴; effnet_B and zoobot_N stay frozen as random-FPR / diversity anchors.",
        "Storfer and Inchausti held out as both positives and negatives ⇒ the gate "
        "measures generalisation, not memorisation.",
    ], size=14)
    add_picture_fit(slide, FIG / "fig3_v4_finetune_gate.png",
                    Inches(2.6), Inches(3.05), Inches(8.1), Inches(3.4))
    add_takeaway(slide, "One short DR11-native fine-tune cracks the lrg+companion hard "
                        "residual v3 could not touch (0.544 → 0.796), with no AUC regression.")


def slide_14_resweep(prs):
    slide = new_slide(prs, notes=(
        "The full v4 re-sweep. All 5.38e7 = 53,809,040 DR11-south galaxies re-scored with "
        "the fine-tuned ensemble, reusing the original parent cutouts (no re-extraction). "
        "Combiner = mean of [effnet_B, zoobot_N, and the three _b50_dr11 members]. Select "
        "the top 150k by mean.\n\n"
        "Held-out recall, position-crossmatch, Storfer/Inchausti training-excluded -- the "
        "best of any configuration (dr11-campaign-v4/papers/main.tex Table 2):\n"
        "  Inchausti grade-A: union-95k 0.54 | mean-150k 0.80 | v4 0.87 (60/69)\n"
        "  Inchausti grade-B: v4 0.87 (76/87)\n"
        "  Storfer grade-A (hard): union-95k 0.32 | mean-150k 0.61 | v4 0.825 (85/103)\n\n"
        "The full re-sweep recovers out-of-pool lenses that a survivor-only re-rank cannot: "
        "the v4 top-150k retains only 15,922 of the union-95k survivors and ADDS 134,078.\n\n"
        "STATUS: this candidate set is recall-richer but PENDING VETTING. The completed, "
        "HSC-vetted catalogue is still the v3blend8 as-run set. Sections 4-5 of the "
        "dr11-campaign-v4 report (candidate catalogue + LensJudge HSC vetting, conclusion) "
        "are a scaffold, marked TODO.\n\n"
        "If asked why mean-150k beats union-95k by so much: it is not only the extra 55k "
        "budget -- the union is dominated by single-member spikes, see slide 12."
        f"\n\n{CITE}"))
    add_title(slide, "v4 — the full re-sweep and the candidate set")
    add_subtitle(slide, "All 53,809,040 DR11-south galaxies re-scored with the fine-tuned "
                        "ensemble (original parent cutouts reused); top 150k by mean")
    add_picture_fit(slide, FIG / "fig4_selector_bottleneck.png",
                    Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.35))
    box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.05), CONTENT_W, Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    set_text(tf, "Status: the v4 candidate set is recall-richer but PENDING VETTING. "
                 "The completed, HSC-vetted catalogue remains the v3blend8 as-run set.",
             size=14, bold=True, color=ACCENT)
    add_takeaway(slide, "Best recall of any configuration by a wide margin — and 89% of "
                        "the top-150k is territory the union sweep never surfaced.",
                 color=NAVY)


def slide_15_tension(prs):
    slide = new_slide(prs, notes=(
        "The central tension. Keep the two columns rigidly apart.\n\n"
        "LEFT (established): v2-lean over v1 at 0.01% FPR: Storfer +0.340 [+0.306,+0.381], "
        "Inchausti +0.256 [+0.217,+0.301], paired bootstrap 10,000 reps over 1M negatives. "
        "v4 re-sweep: Inchausti-A 0.54->0.87, Storfer-A 0.32->0.825.\n\n"
        "RIGHT (bounded): the v3 DR10 sweep produced ZERO net-new lenses -- the four 'new' "
        "grade-A survivors are all previously published (Huang+2020 DESI-038.2078-03.3906; "
        "CSWA 1, the Cosmic Horseshoe; DES J0352-3825; AGEL J221912-434835). v3's headline "
        "49 cross-validated A/B are RECOVERIES of published Euclid-Q1 lenses, not "
        "discoveries. The DR11 campaign yielded 3 genuinely new after HSC vetting.\n\n"
        "Why: 'no net-new lens is established from DECaLS pixels alone.' The resolution "
        "lever CONVERTS overlap candidates; it does not manufacture discoveries.\n\n"
        "The nuance to land: v4's improvements are real and large on the metric we can "
        "measure. They have NOT yet been shown to convert into new confirmed lenses."))
    add_title(slide, "What the DR11 sweep discovered — and the tension to hold")

    lbl = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.2), Inches(6.0), Inches(0.4))
    set_text(lbl.text_frame, "Recall / recovery — large, CI'd, established",
             size=16, bold=True, color=GOOD)
    body(slide, 1.68, 4.6, [
        "v2-lean over v1 @0.01% FPR (paired bootstrap, 10,000 reps, 1M negatives):",
        (1, "Storfer  +0.340  [+0.306, +0.381]"),
        (1, "Inchausti  +0.256  [+0.217, +0.301]"),
        "v4 re-sweep, held out and training-excluded:",
        (1, "Inchausti grade-A  0.54 → 0.87"),
        (1, "Storfer grade-A (hard)  0.32 → 0.825"),
        "v4 fine-tune gate passes on all three held-out sets, AUC 0.996 → 0.998.",
        "Every one of these is a measured, interval-bearing number.",
    ], size=14, left_in=0.5, width_in=6.0)

    lbl2 = slide.shapes.add_textbox(Inches(6.95), Inches(1.2), Inches(5.9), Inches(0.4))
    set_text(lbl2.text_frame, "Discovery — small, resolution-bounded",
             size=16, bold=True, color=ACCENT)
    body(slide, 1.68, 4.6, [
        "v3 DR10 sweep: zero net-new lenses. The four \"new\" grade-A survivors are all "
        "previously published — Huang+2020 DESI-038.2078−03.3906, CSWA 1 (the Cosmic "
        "Horseshoe), DES J0352−3825, AGEL J221912−434835.",
        "v3's headline 49 cross-validated A/B are recoveries of published Euclid-Q1 "
        "lenses, not discoveries.",
        "v3 DR11 campaign: 3 genuinely new systems after HSC tier-2 vetting.",
        "v4: 134,078 new candidates, none vetted yet.",
        "\"No net-new lens is established from DECaLS pixels alone.\" The resolution "
        "lever converts overlap candidates; it does not manufacture discoveries.",
    ], size=14, left_in=6.95, width_in=5.9)

    add_takeaway(slide, "We got dramatically better at ranking lenses and only marginally "
                        "better at confirming new ones — because the two are limited by "
                        "different things: the combiner and the weights vs. the pixel scale.")


def slide_16_confidence_i(prs):
    slide = new_slide(prs, notes=(
        "Confidence, part I. The hardest, most independent test says the models tie.\n\n"
        "Independent Euclid grade-A-vs-C AUC (claudenet/data/v3/model_progression.json, "
        "v3_section.tex tab_progression):\n"
        "  orig (effnet_B, random-neg) 0.613 | v2-lean 0.647 | v3blend8 0.660 | "
        "v3-head 0.615.  SE ~= 0.06, n_A = 66.\n"
        "All four sit within about one standard error -- statistically indistinguishable "
        "at the hardest distinction (real lens vs convincing mimic at Euclid resolution).\n\n"
        "RETRACTED: the seed mimic-FPR metric is selection-biased -- the seed bank IS what "
        "the v2/v3 ensemble scored high -- so any 'Nx over the original published models' "
        "claim read off it is invalid, and we do not make it.\n\n"
        "What IS solid: v3 clearly beats the original CNN on the broad mimic population "
        "(DR10-eval recovery 0.598 -> 0.885); and LensJudge escalation flips median p_lens "
        "from 0.03 to 0.70 at 0.1 arcsec on 134 real objects, with ~90% agreement with "
        "Euclid experts on grade-A.\n\n"
        "Conclusion: the discovery frontier is resolution-bounded for every model. "
        "Architecture and training gains do not resolve the lens-vs-mimic tie at Euclid "
        "resolution."))
    add_title(slide, "How confident are we? (I) The hardest test says the models tie")
    add_picture_fit(slide, FIG / "fig5_confidence_frontier.png",
                    Inches(0.45), Inches(1.2), Inches(8.05), Inches(3.6))
    add_picture_fit(slide, CN_FIG / "euclid_flip.png",
                    Inches(8.75), Inches(1.2), Inches(4.15), Inches(3.6))
    body(slide, 4.95, 1.55, [
        "Retracted: the seed mimic-FPR metric is selection-biased — the seed bank is "
        "precisely what the v2/v3 ensemble scored high — so any \"N× over the original "
        "published models\" claim read off it is invalid. We do not make it.",
        "Still solid: broad-mimic recovery 0.598 → 0.885 across the four generations "
        "(right panel), and the LensJudge escalation flip (median p_lens 0.03 → 0.70 "
        "at 0.1″ on 134 real objects, ~90% agreement with Euclid experts on grade-A).",
    ], size=13)
    add_takeaway(slide, "On the one truly independent test, all four generations tie within "
                        "a standard error. Resolution breaks the lens-vs-mimic tie; models don't.")


def slide_17_confidence_ii(prs):
    slide = new_slide(prs, notes=(
        "Confidence, part II. A ledger. Be precise about tense.\n\n"
        "SOLID and interval-bearing: the matched-FPR recovery gains (v2 paired bootstrap "
        "over 1M negatives; v4 gate and re-sweep on held-out, training-excluded catalogs).\n"
        "SOLID and vetted: 3 genuinely new lenses + 39 new DESI-resolution-only A/B from "
        "the v3 DR11 campaign, HSC tier-2 confirmed.\n"
        "PENDING: the whole v4 candidate set. 134,078 new candidates, zero vetted.\n"
        "RETRACTED: the seed mimic-FPR 'Nx' claim.\n"
        "POWER-LIMITED: conformal selection certifies nothing at DR11 scale -- a surfaced "
        "power floor, not a null result.\n\n"
        "And the provenance point, if it comes up: the v4 result artifacts live on "
        "Perlmutter scratch, not in the repo. Every v4 number in this deck is transcribed "
        "from v4_section.tex and dr11-campaign-v4/papers/main.tex. The v1/v2/v3 numbers "
        "are backed by tracked JSON under claudenet/data/."))
    add_title(slide, "How confident are we? (II) Solid vs pending")
    add_subtitle(slide, "Tense discipline: what is measured, what is vetted, what is "
                        "merely promising")
    add_table(slide, Inches(1.62), Inches(4.3), [
        ["claim", "evidence", "interval", "vetted?", "status"],
        ["v2-lean ≫ v1 ≫ published meta at matched FPR",
         "NegEval-1M, 1M negatives", "paired bootstrap 95% CI", "n/a", "SOLID"],
        ["v4 fine-tune gate: +0.09 / +0.30 / +0.25",
         "held-out, training-excluded", "no CI reported", "n/a", "SOLID"],
        ["v4 re-sweep recall 0.87 / 0.825",
         "position-crossmatch, held out", "no CI reported", "n/a", "SOLID"],
        ["3 genuinely new lens systems",
         "LensJudge + HSC tier-2 + literature", "—", "yes", "CONFIRMED"],
        ["v4's 134,078 new candidates",
         "stage-1 mean only", "—", "no", "PENDING"],
        ["\"N× better than the published models\"",
         "seed mimic-FPR (selection-biased)", "—", "—", "RETRACTED"],
        ["conformal-certified selection at DR11 scale",
         "group-conformal BH", "power floor reached", "—", "CERTIFIES NOTHING"],
    ], col_widths_in=[3.7, 3.1, 2.35, 1.05, 2.13], size=11)
    add_takeaway(slide, "The recall gains are established and v3's confirmed lenses are "
                        "real. v4's new pool is promising but unvetted — so we make no "
                        "discovery claim for it today.")


def slide_18_scope(prs):
    slide = new_slide(prs, notes=(
        "Honest scope. Two structural limits remain, and neither is fixed by a better "
        "model.\n\n"
        "1. v4 remains DECam-specific. The northern footprint (BASS/MzLS) still "
        "under-scores systematically -- the v3 DR11-north sweep had zero Euclid-A/B in its "
        "top-30. It needs north-native retraining.\n"
        "2. The DECaLS 0.262\"/px resolution ceiling is unchanged. Higher-resolution "
        "vetting remains the decisive net-new-discovery lever. About 40% of Euclid "
        "grade-A/B lenses are not even in the DESI parent -- too faint or too small. "
        "No model can cross that.\n\n"
        "Robustness we did check: the legacy Gaussian high-pass residual was replaced by a "
        "signed chi = (data - model)/sigma, and the entire 500-candidate cascade was re-run "
        "in a paired design. The new residual runs ~2x hotter at tier-1 (raw A/B 90 -> 135) "
        "but the HSC tier-2 gate absorbs it (tier-2 A/B 24 -> 27). All 9 new-class systems "
        "survive. Conclusions unchanged.\n\n"
        "Figures: fig7 (this deck) and dr11-campaign/papers/figures/residual_robustness.png."))
    add_title(slide, "Honest scope — what ClaudeNet fixed, and what it cannot")
    body(slide, 1.18, 1.5, [
        "v4 remains DECam-specific: the BASS/MzLS north still under-scores "
        "systematically (v3's DR11-north top-30 contained zero Euclid A/B).",
        "The 0.262″/px ceiling is unchanged, and ~40% of Euclid grade-A/B lenses "
        "are not even in the DESI parent — too faint or too small for any DESI model.",
        "Robustness holds: replacing the legacy residual with a signed χ=(data−model)/σ "
        "runs ~2× hotter at tier-1 (raw A/B 90→135), but the HSC tier-2 gate absorbs it "
        "(24→27) and all 9 new-class systems survive.",
    ], size=13)
    add_picture_fit(slide, FIG / "fig7_bottlenecks_fixed.png",
                    Inches(0.5), Inches(2.85), Inches(6.1), Inches(3.55))
    add_picture_fit(slide, D11_FIG / "residual_robustness.png",
                    Inches(6.85), Inches(2.85), Inches(5.95), Inches(3.55))
    add_takeaway(slide, "Two structural limits remain — the northern instrument and the "
                        "DECam pixel scale — and neither is fixed by a better model.")


def slide_19_future(prs):
    slide = new_slide(prs, notes=(
        "Future work, verbatim priorities from claudenet/papers/future_work.tex.\n\n"
        "1. Euclid pre-screening -- the decisive discovery lever. v3/v4 pre-screen DESI; "
        "Euclid confirms at 0.1 arcsec. Q1's ~63 deg^2 over three Deep Fields overlaps "
        "only ~0.5% of a DESI-south sweep, which is exactly why net-new yield there is "
        "bounded. Euclid DR1 (wide-area) inverts this. The ready action is a targeted "
        "sweep restricted to the Euclid-overlap sky, where every candidate is escalatable.\n"
        "2. BASS/MzLS-north retraining. Mine a north-native lens-mimic bank and retrain the "
        "members -- the recipe is instrument-agnostic, only the negatives change. Roughly "
        "triples the addressable footprint.\n"
        "3. i-band / foundation-model revival for small theta_E. Re-test the deferred "
        "AION/i-band lever under the NARROWER lens-vs-mimic task, not the random-FPR task "
        "where it failed in v2. Targets the red-lens/red-companion degeneracy.\n"
        "4. A pixel-level lens-vs-mimic head. The learned logistic head on member scores "
        "overfits the strict tail and does not generalise to Euclid; train a small CNN head "
        "directly on pixels (positives=lenses, negatives=the typed mimic bank) as a "
        "post-conformal re-ranker.\n"
        "5. The active-learning loop. LensJudge exports confident-D rejections -- true "
        "mimics, NOT objects that flipped to A/B at higher resolution -- as fresh typed "
        "hard negatives, closing a mine->grade->retrain loop.\n"
        "6. A full DR11-south candidate catalogue: realised, that is v4."))
    add_title(slide, "Future work — the levers, in priority order")
    body(slide, 1.2, 5.15, [
        "1.  Euclid pre-screening — the decisive discovery lever. v4 pre-screens DESI; "
        "Euclid confirms at 0.1″.",
        (1, "Q1's ~63 deg² over three Deep Fields overlaps only ~0.5% of a DESI-south "
            "sweep — which is exactly why net-new yield is bounded there. Euclid DR1 "
            "(wide-area) inverts this."),
        (1, "Ready action: a targeted sweep restricted to the Euclid-overlap sky, where "
            "every candidate is escalatable."),
        "2.  BASS/MzLS-north retraining — mine a north-native mimic bank and retrain. "
        "The recipe is instrument-agnostic; only the negatives change. ~3× the "
        "addressable footprint.",
        "3.  i-band / foundation-model revival for small θ_E — re-test the deferred "
        "AION/i-band lever under the narrower lens-vs-mimic task, not the random-FPR "
        "task where it failed in v2.",
        "4.  A pixel-level lens-vs-mimic head — the logistic head on member scores "
        "overfits the strict tail and doesn't generalise to Euclid; train a small CNN "
        "on pixels as a post-conformal re-ranker.",
        "5.  The active-learning loop — LensJudge exports confident-D rejections (true "
        "mimics, not resolution-flips) as fresh typed hard negatives: mine → grade → retrain.",
    ], size=14, width_in=7.4)
    add_picture_fit(slide, CN_FIG / "euclid_recall.png",
                    Inches(8.2), Inches(1.9), Inches(4.7), Inches(3.6))
    add_takeaway(slide, "The next real discoveries come from where we vet, not from what "
                        "we train. Euclid-overlap sweeps and north-native retraining are "
                        "the two highest-leverage moves.")


def slide_20_next_step(prs):
    slide = new_slide(prs, notes=(
        "The immediate next step and the deployment recommendation.\n\n"
        "IMMEDIATE: run the LensJudge v3 HSC tier-2 cascade on the v4 NEW shortlist "
        "(134,078 new candidates). That fills in sections 4-5 of the dr11-campaign-v4 "
        "report, which are currently a TODO scaffold, and it is the one action that could "
        "convert v4's measured recall gain into an actual vetted yield. The v3 campaign "
        "cost $40.38 for 500 candidates, so budget accordingly for the shortlist.\n\n"
        "DEPLOYMENT (future_work.tex): for DR11-south ship v4 -- the release-native "
        "fine-tuned ensemble with the calibrated-mean stage-1 selector. Retain v3blend8 "
        "for DR10-south and the Euclid cross-validation. In all cases use the score MEAN, "
        "not the learned head (it does not generalise past phi=0.05), plus the LensJudge "
        "cascade for vetting.\n\n"
        "What this buys: keeps v2-lean's lens recall, adds v3's contaminant-rejection "
        "edge, adds v4's release-native gain on the hard residual, and routes the "
        "resolution-breaking confirmation to exactly the candidates that have "
        "higher-resolution coverage."))
    add_title(slide, "Next step, and what to deploy")

    lbl = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.2), CONTENT_W, Inches(0.4))
    set_text(lbl.text_frame, "Immediate", size=17, bold=True, color=ACCENT)
    body(slide, 1.68, 1.5, [
        "Run the LensJudge v3 HSC tier-2 cascade on the v4 NEW shortlist (134,078 "
        "candidates). This fills sections 4–5 of the dr11-campaign-v4 report, which are "
        "still a TODO scaffold.",
        "It is the one action that could convert v4's measured recall gain into an "
        "actual vetted yield. Reference cost: the v3 campaign vetted 500 candidates "
        "for $40.38.",
    ], size=15)

    lbl2 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.35), CONTENT_W, Inches(0.4))
    set_text(lbl2.text_frame, "Deployment recommendation", size=17, bold=True, color=NAVY)
    body(slide, 3.83, 2.6, [
        "DR11-south: ship v4 — the release-native fine-tuned ensemble with the "
        "calibrated-mean stage-1 selector.",
        "DR10-south and the Euclid cross-validation: retain v3blend8.",
        "In all cases combine with the score mean, not the learned head — it does not "
        "generalise past φ = 0.05 — and vet with the LensJudge cascade.",
        "This keeps v2-lean's lens recall, adds v3's contaminant rejection and v4's "
        "release-native gain on the hard residual, and routes the resolution-breaking "
        "confirmation to exactly the candidates with higher-resolution coverage.",
    ], size=15)
    add_takeaway(slide, "Ship v4 with the mean selector for DR11-south now; vet its 134k "
                        "new candidates next. That is the step that could turn a recall "
                        "gain into new lenses.")


def slide_21_summary(prs):
    slide = new_slide(prs, notes=(
        "Summary. Land the two-part claim cleanly and stop.\n\n"
        "We built a demonstrably better lens-ranker, and we drew a clean line between what "
        "that proves (recall) and what it does not yet prove (discovery).\n\n"
        "The path forward is resolution, not architecture.\n\n"
        "Cutouts: the same three genuinely new systems as the title slide "
        "(s_327526_9059, s_340971_2860, s_325163_2981), DESI grz over HSC grizy."))
    add_title(slide, "Summary")
    body(slide, 1.2, 3.05, [
        "The lineage's ceiling was never architecture — a 194k ResNet ties a 20.5M "
        "EfficientNet, and the published meta-learner was a simple average in disguise.",
        "ClaudeNet's levers, in the order each exposed the next: diversity (v1) → honest "
        "deployment-scale negatives (v2) → contaminant awareness (v3) → release-native "
        "adaptation (v4).",
        "v4 = v3's members + a calibrated-mean selector (free) + a DR11-native fine-tune "
        "(one GPU pass). Best recall of any configuration: Inchausti grade-A 0.54 → 0.87.",
        "The honest split: recall gains are large and CI'd; net-new discovery is three "
        "confirmed systems from the v3 DR11 campaign, and v4's 134,078 new candidates "
        "remain unvetted.",
        "On the one independent test, all four generations tie within one standard error. "
        "The frontier is resolution-bounded.",
    ], size=15)

    cap = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.35), CONTENT_W, Inches(0.35))
    set_text(cap.text_frame,
             "3 confirmed  ·  134,078 pending — the three genuinely new systems "
             "(DESI grz, HSC grizy)", size=12, color=MUTED)
    for i, sysid in enumerate(NEW3):
        x = Inches(1.35 + i * 3.9)
        add_picture_fit(slide, D11_FIG / f"cand_s_{sysid}_desi.png",
                        x, Inches(4.75), Inches(1.4), Inches(1.4))
        add_picture_fit(slide, D11_FIG / f"cand_s_{sysid}_hsc.png",
                        Inches(2.85 + i * 3.9), Inches(4.75), Inches(1.4), Inches(1.4))
        lab = slide.shapes.add_textbox(Inches(1.35 + i * 3.9), Inches(6.2),
                                       Inches(2.9), Inches(0.35))
        set_text(lab.text_frame, f"s_{sysid}", size=11, color=INK)

    add_takeaway(slide, "We built a demonstrably better lens-ranker — and drew a clean line "
                        "between what that proves (recall) and what it does not yet prove "
                        "(discovery). The path forward is resolution, not architecture.",
                 top=Inches(6.68))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for fn in (
        slide_01_cover, slide_02_bluf, slide_03_lineage, slide_04_negative_result,
        slide_05_metric, slide_06_v1_diversity, slide_07_v2, slide_08_v2_reality,
        slide_09_v3, slide_10_v3_campaign, slide_11_v4_overview, slide_12_lever1,
        slide_13_lever2, slide_14_resweep, slide_15_tension, slide_16_confidence_i,
        slide_17_confidence_ii, slide_18_scope, slide_19_future, slide_20_next_step,
        slide_21_summary,
    ):
        fn(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}  ({OUTPUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
