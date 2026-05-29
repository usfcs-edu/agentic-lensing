#!/usr/bin/env python3
"""
Build the progress deck for the Huang-group strong-lensing reproduction sprint.

Wraps python-pptx to construct a ~16-slide deck summarizing five reproductions
(Huang 2020, Huang 2021 [in progress], Foundry I, Hsu 2025, and the adjacent
SpectrumFM/redshifty foundation-model track) for the research group meeting.

Style (palette, helpers, 16:9 blank layout) mirrors
tools/spectrumfm/build_deck.py so the deck looks native to the group's decks.

Run with the venv that has python-pptx (1.0.2):
    /raid/benson/.venvs/redshifty/bin/python3 progress/2026-05-29/build_deck.py
Output:
    progress/2026-05-29/huang_reproductions_2026-05-29.pptx
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO_ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = REPO_ROOT / "progress" / "2026-05-29"
OUTPUT = OUT_DIR / "huang_reproductions_2026-05-29.pptx"

# Figure assets (verified on disk)
F_H20_ARCH = REPO_ROOT / "reproductions/huang-2020/papers/figures/lanusse_resnet_arch.png"
F_H20_ROC = REPO_ROOT / "reproductions/huang-2020/papers/figures/roc_curve.png"
F_H20_REC = REPO_ROOT / "reproductions/huang-2020/papers/figures/recovery_by_grade.png"
F_H21_ARCH = REPO_ROOT / "reproductions/huang-2021/papers/figures/shielded_arch.png"
F_H21_ROC = REPO_ROOT / "reproductions/huang-2021/papers/figures/shielded_roc.png"
F_H21_CMP = REPO_ROOT / "reproductions/huang-2021/papers/figures/arch_comparison.png"
F_FI_CUTOUT = REPO_ROOT / "reproductions/foundry-i/figs/cutout_preview.png"
F_FI_RESID = REPO_ROOT / "reproductions/foundry-i/figs/map_residual.png"
F_FI_PSF = REPO_ROOT / "reproductions/foundry-i/figs/psf_comparison.png"
F_HSU_TBL = REPO_ROOT / "reproductions/hsu-2025/figs/inspect_table2.jpg"
F_RS_DIV = REPO_ROOT / "reproductions/redshifty/papers/figures/diversity_journey.png"

# 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
LEFT_MARGIN = Inches(0.5)
TOP_MARGIN = Inches(0.4)
CONTENT_W = SLIDE_W - 2 * LEFT_MARGIN

NAVY = RGBColor(0x18, 0x2C, 0x5B)
ACCENT = RGBColor(0xC6, 0x3A, 0x3A)
GOOD = RGBColor(0x1E, 0x7A, 0x3C)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)


def set_text(textframe, text, size=24, bold=False, color=INK, align=None):
    textframe.clear()
    p = textframe.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, text, size=30):
    box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, CONTENT_W, Inches(0.8))
    set_text(box.text_frame, text, size=size, bold=True, color=NAVY)
    return box


def add_subtitle(slide, text, top=Inches(1.1), size=18, color=MUTED, width=None):
    box = slide.shapes.add_textbox(LEFT_MARGIN, top, width or CONTENT_W, Inches(0.5))
    box.text_frame.word_wrap = True
    set_text(box.text_frame, text, size=size, color=color)
    return box


def add_body(slide, top, height, bullets, size=18, indent_size=14, width=None, left=None):
    """bullets: list[str] or list[tuple(level, str)]."""
    box = slide.shapes.add_textbox(left or LEFT_MARGIN, top, width or CONTENT_W, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = ("• " if level == 0 else "– ") + text
        run.font.size = Pt(size if level == 0 else indent_size)
        run.font.color.rgb = INK
    return box


def add_quote(slide, top, height, text, size=18, width=None):
    box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.4), top,
                                   (width or CONTENT_W) - Inches(0.8), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "“" + text + "”"
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.color.rgb = NAVY
    return box


def add_footer(slide, text):
    box = slide.shapes.add_textbox(LEFT_MARGIN, SLIDE_H - Inches(0.45),
                                   CONTENT_W, Inches(0.3))
    set_text(box.text_frame, text, size=10, color=MUTED)


def add_badge(slide, text, color, left, top, width=Inches(2.2)):
    """A small filled status pill."""
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.42))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shp


def add_table(slide, top, height, rows, col_widths_in, header_row=True, size=14, left=None):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        left or LEFT_MARGIN, top,
        Inches(sum(col_widths_in)), height,
    )
    table = table_shape.table
    for col_idx, w in enumerate(col_widths_in):
        table.columns[col_idx].width = Inches(w)
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(size)
            if r_idx == 0 and header_row:
                run.font.bold = True
                run.font.color.rgb = NAVY
            else:
                run.font.color.rgb = INK
    return table_shape


def add_picture_fit(slide, path, left, top, max_w, max_h, caption=None):
    """Embed an image scaled to fit (left,top,max_w,max_h), centered. Falls
    back to a captioned placeholder box if the file is missing/unreadable."""
    path = Path(path)
    try:
        from PIL import Image
        with Image.open(path) as im:
            iw, ih = im.size
        aspect = iw / ih
        box_aspect = max_w / max_h
        if aspect >= box_aspect:
            draw_w = max_w
            draw_h = int(max_w / aspect)
        else:
            draw_h = max_h
            draw_w = int(max_h * aspect)
        off_l = left + (max_w - draw_w) // 2
        off_t = top + (max_h - draw_h) // 2
        slide.shapes.add_picture(str(path), off_l, off_t, width=draw_w, height=draw_h)
        if caption:
            cap = slide.shapes.add_textbox(left, top + max_h, max_w, Inches(0.3))
            set_text(cap.text_frame, caption, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    except Exception as exc:  # noqa: BLE001 - never let a missing fig break the build
        ph = slide.shapes.add_textbox(left, top, max_w, max_h)
        tf = ph.text_frame
        tf.word_wrap = True
        set_text(tf, f"[figure unavailable: {path.name}]\n{exc}", size=12, color=MUTED,
                 align=PP_ALIGN.CENTER)


# ----------------------------------------------------------------------------- slides

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.3), CONTENT_W, Inches(1.6))
    set_text(box.text_frame,
             "Reproducing the Huang-Group\nStrong-Lensing Program",
             size=40, bold=True, color=NAVY)
    box2 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.2), CONTENT_W, Inches(0.6))
    set_text(box2.text_frame,
             "A one-week onboarding sprint — five results re-derived from the papers",
             size=21, color=ACCENT)
    box3 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(5.2), CONTENT_W, Inches(0.4))
    set_text(box3.text_frame,
             "Greg Benson  ·  University of San Francisco  ·  29 May 2026  ·  for the Huang Lensing Group",
             size=16, color=INK)


def slide_program(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The program — and what this sprint covered")
    add_subtitle(slide,
                 "Huang group: ML discovery of strong lenses in DESI surveys, then follow-up + modeling. "
                 "~5,500 candidates across DR8/9/10 — an order-of-magnitude expansion of the field.",
                 top=Inches(1.05), size=15)
    add_table(slide, Inches(1.9), Inches(2.5), [
        ["Pillar", "What it does", "Reproduced here"],
        ["Discover", "Five complementary modalities find lenses in DESI imaging + spectra",
         "Huang 2020 & 2021 (image finders); Hsu 2025 (spectroscopic pairs)"],
        ["Characterize", "DESI Strong Lens Foundry — HST/DESI/Keck/VLT follow-up + GIGA-Lens modeling",
         "Foundry I (GIGA-Lens fit of an HST-confirmed system)"],
        ["Cosmologize", "H₀ and dark-matter substructure (downstream science)", "— (future)"],
    ], col_widths_in=[1.9, 5.6, 4.8], size=14)
    add_subtitle(slide,
                 "This sprint reproduced 5 results across these pillars — 4 of them with NO public code "
                 "(re-implemented from the methods sections) — on local GPUs, in ~1 week.",
                 top=Inches(5.0), size=16, color=ACCENT)
    add_subtitle(slide,
                 "Plus an adjacent foundation-model track: the SpectrumFM / redshifty redshift-ignition result.",
                 top=Inches(5.7), size=14, color=MUTED)
    add_footer(slide, "Onboarding context: plans/AGENTIC_LENSING_ONBOARDING_PLAN.md (16-paper landscape).")


def slide_scoreboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Scoreboard — five reproductions")
    add_table(slide, Inches(1.3), Inches(4.2), [
        ["Effort", "Pillar", "Reproduced result", "vs published", "Status"],
        ["Huang 2020", "Discover",
         "Test AUC 0.9943 (DR7) / 0.9991 (DR9)", "paper val AUC 0.98", "✓ complete"],
        ["Huang 2021", "Discover",
         "Shielded net 59,905 params (58.6× smaller); AUC ≈ L18 within ±0.002", "matches central claim",
         "in progress (DR8)"],
        ["Foundry I", "Characterize",
         "θ_E to 3.0%; all 6 mass-parameter signs recovered", "GIGA-Lens on HST", "✓ complete"],
        ["Hsu 2025", "Discover",
         "13,530 pairs / 27,334 spectra; 100% recall on Table 2", "paper 13,218 / 26,621", "✓ complete"],
        ["SpectrumFM", "Foundation",
         "Redshift ignition; AR/TF ratio 0.73", "NERSC ref 0.74", "✓ reproduced"],
    ], col_widths_in=[1.7, 1.6, 5.3, 2.4, 1.6], size=13)
    add_subtitle(slide,
                 "Detail on each follows. Absolute finder AUCs carry a known training-leakage caveat (see integrity slide).",
                 top=Inches(6.4), size=13, color=MUTED)


def slide_h20_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Huang 2020 — image lens-finder (Discover)")
    add_subtitle(slide, "DESI DECaLS ResNet classifier — re-implemented from scratch (no public code)",
                 top=Inches(1.05), size=16, width=Inches(7.7))
    add_body(slide, Inches(1.7), Inches(4.5), [
        "Architecture: Lanusse 2018 / CMU-DeepLens ResNet-46 (3.5M params), reimplemented in PyTorch from Lanusse §3.1.",
        "Inputs: 101×101 px grz cutouts (0.262″/px); 949 NeuraLens L18 lens positives + 5,000 random DR1 negatives.",
        "Training: 120 epochs, Adam, stratified 70/20/10 split.",
        (1, "~25 minutes on a single L4 GPU."),
        "Deployed full-sky via a brick-driven inference pipeline (see results).",
    ], size=16, width=Inches(7.6))
    add_picture_fit(slide, F_H20_ARCH, Inches(8.3), Inches(1.5), Inches(4.6), Inches(4.6),
                    caption="Lanusse / CMU-DeepLens ResNet-46")
    add_footer(slide, "reproductions/huang-2020/ — scripts 01–16 + papers/main.pdf")


def slide_h20_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Huang 2020 — results")
    add_table(slide, Inches(1.3), Inches(2.0), [
        ["Training set", "Val AUC", "Test AUC"],
        ["DR9", "0.9983", "0.9991"],
        ["DR7 (paper-exact)", "0.9890", "0.9943"],
        ["Paper reported", "0.98", "—"],
    ], col_widths_in=[3.0, 2.2, 2.2], size=15, left=Inches(0.5))
    add_body(slide, Inches(3.7), Inches(2.8), [
        "Recovery of published candidates @ p ≥ 0.9: Grade A 83%, B 64%, C 28%.",
        "Brick-driven inference: ~200× faster than the per-object cutout endpoint",
        (1, "full DR7 footprint in ~10 hours instead of ~84 days."),
    ], size=16, width=Inches(7.0))
    add_picture_fit(slide, F_H20_ROC, Inches(8.0), Inches(1.4), Inches(4.9), Inches(4.6),
                    caption="ROC — held-out test set (DR9-trained, AUC 0.9991)")
    add_footer(slide, "Test AUC reaches 0.99+ vs the paper's 0.98 — caveat on absolute AUC discussed on the integrity slide.")


def slide_h21_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Huang 2021 — shielded ResNet")
    add_badge(slide, "IN PROGRESS", ACCENT, Inches(8.9), Inches(0.42), width=Inches(2.3))
    add_subtitle(slide, "The paper's headline innovation: a tiny “shielded” net that matches the big one",
                 top=Inches(1.05), size=16, width=Inches(7.7))
    add_body(slide, Inches(1.7), Inches(4.2), [
        "Shield = 1×1 convolution layers inserted between every three residual blocks (4 shields; final block 32 channels).",
        "Controlled comparison: same DR9/DR7 cutouts, positives, negatives, and split as Huang 2020 — architecture is the ONLY variable.",
        "Shielded net: 59,905 parameters = 58.6× smaller than the 3.5M-parameter L18 baseline.",
    ], size=16, width=Inches(7.6))
    add_picture_fit(slide, F_H21_ARCH, Inches(8.3), Inches(1.5), Inches(4.6), Inches(4.6),
                    caption="Shielded architecture schematic")
    add_footer(slide, "reproductions/huang-2021/ — Phase 4a complete; DR8 deployment (4b/4c) in progress.")


def slide_h21_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Huang 2021 — shielded matches (and beats) L18")
    add_table(slide, Inches(1.25), Inches(2.2), [
        ["Run", "Params", "Val AUC", "Test AUC"],
        ["L18 / DR9", "3,508,833", "0.9983", "0.9991"],
        ["Shielded / DR9", "59,905", "0.9989", "0.9988"],
        ["L18 / north-aug", "3,508,833", "0.9991", "0.9985"],
        ["Shielded / north-aug", "59,905", "0.9992", "0.9996"],
    ], col_widths_in=[2.7, 1.9, 1.6, 1.6], size=13, left=Inches(0.5))
    add_subtitle(slide,
                 "Shielded matches L18 within ±0.002 AUC at 59× fewer params — and edges ahead on "
                 "north-augmented data (0.9996 vs 0.9985).",
                 top=Inches(3.7), size=15, color=ACCENT, width=Inches(7.2))
    add_body(slide, Inches(4.7), Inches(2.3), [
        "DR8 full-sky deployment (in progress): two-model (L18 + shielded) brick inference over ~23M parent galaxies "
        "across South (DECaLS) + North (BASS/MzLS) footprints.",
        (1, "All 4 shards scored; Phase 4c recovery + leakage crossmatch pending."),
    ], size=14, width=Inches(7.2))
    add_picture_fit(slide, F_H21_ROC, Inches(8.0), Inches(1.4), Inches(4.9), Inches(4.6),
                    caption="ROC — shielded vs L18")


def slide_fi_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Foundry I — GIGA-Lens Bayesian modeling (Characterize)")
    add_subtitle(slide, "Differentiable lens modeling of an HST-confirmed Foundry system",
                 top=Inches(1.05), size=16, width=Inches(7.7))
    add_body(slide, Inches(1.7), Inches(4.3), [
        "Target: DESI-165.4754−06.0423, HST WFC3/F140W (GO-15867, public MAST).",
        "Model: 74 parameters — EPL + external shear + 4 Sérsic lens-light + Sérsic source + shapelets.",
        "Inference: 200-chain multi-start MAP + 200-particle SVI (+ NUTS diagnostic) on 2× L4, ~7 min.",
        "Empirical PSF built from in-field stars — critical for any GIGA-Lens reproduction.",
    ], size=16, width=Inches(7.6))
    add_picture_fit(slide, F_FI_CUTOUT, Inches(8.3), Inches(1.5), Inches(4.6), Inches(4.6),
                    caption="HST WFC3/F140W cutout")
    add_footer(slide, "reproductions/foundry-i/ — scripts 01–25 + papers/main.pdf")


def slide_fi_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Foundry I — posterior vs Huang 2025a")
    add_table(slide, Inches(1.3), Inches(2.4), [
        ["Parameter", "Reproduced (v10 SVI)", "Published (HMC)", "Match"],
        ["θ_E (″)", "2.566", "2.646", "3.0%"],
        ["e₁", "0.106", "0.109", "2.5%"],
        ["shear position angle", "—", "—", "within 1°"],
        ["γ_EPL (slope)", "2.15", "1.37", "57% steeper (open)"],
    ], col_widths_in=[2.8, 2.5, 2.1, 2.2], size=13, left=Inches(0.5))
    add_subtitle(slide,
                 "Recovers all six mass-parameter sign quadrants; θ_E to 3.0% (NUTS second mode to 1.6%).",
                 top=Inches(4.0), size=15, color=GOOD, width=Inches(7.2))
    add_subtitle(slide,
                 "The γ_EPL slope is 57% steeper than published — flagged as a known open discrepancy "
                 "(likely one missing constraint), not hidden.",
                 top=Inches(4.8), size=14, color=ACCENT, width=Inches(7.2))
    add_picture_fit(slide, F_FI_RESID, Inches(8.0), Inches(1.4), Inches(4.9), Inches(4.6),
                    caption="Best-fit model residual")


def slide_hsu_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Hsu 2025 — spectroscopic lens-pair search (Discover)")
    add_subtitle(slide, "Finding “dimple” lenses in 28M DESI DR1 spectra",
                 top=Inches(1.05), size=16, width=Inches(7.7))
    add_body(slide, Inches(1.7), Inches(4.3), [
        "Pipeline: pre-filter (ZCAT_PRIMARY ∧ ZWARN=0 ∧ ≠STAR ∧ z>0) → spherimatch friends-of-friends "
        "(3″ link) → redshift-ratio cut (z_max/z_min ≥ 1.3) → Einstein-radius classifier.",
        "Uses the paper's own public tool: technic960183/spherimatch.",
        "Runs on 28.4M spectra in ~2 minutes on commodity hardware (FoF step ~36 s).",
    ], size=16, width=Inches(7.6))
    add_picture_fit(slide, F_HSU_TBL, Inches(8.3), Inches(1.5), Inches(4.6), Inches(4.6),
                    caption="Table-2 candidate inspection grid")
    add_footer(slide, "reproductions/hsu-2025/ — scripts 01–10 + papers/main.pdf")


def slide_hsu_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Hsu 2025 — algorithmic match within ~2%")
    add_table(slide, Inches(1.4), Inches(2.6), [
        ["Stage", "Reproduced", "Published"],
        ["After pre-filter", "15,786,243", "~15.8M"],
        ["FoF groups (z-ratio ≥ 1.3)", "13,530", "13,218"],
        ["Spectra in retained groups", "27,334", "26,621"],
    ], col_widths_in=[4.6, 3.2, 3.0], size=16)
    add_subtitle(slide,
                 "Load-bearing counts reproduce to +2.4% / +2.7% — inside the ±5% validation tolerance.",
                 top=Inches(4.4), size=17, color=GOOD)
    add_subtitle(slide,
                 "100% recall on the 20 Grade-A new candidates in Table 2 (median offset 0.91″).",
                 top=Inches(5.1), size=16, color=NAVY)
    add_subtitle(slide,
                 "Visual-inspection grading (the paper's 2,046 / 318 final counts) is out of scope — it needs the team.",
                 top=Inches(5.8), size=14, color=MUTED)
    add_footer(slide, "Public DESI DR1 + spherimatch; fully reproducible on commodity hardware in ~2 minutes.")


def slide_spectrumfm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "SpectrumFM / redshifty — adjacent foundation-model track")
    add_subtitle(slide, "Reproducing the Approach-A “redshift ignition” on commodity GPUs (29 May)",
                 top=Inches(1.05), size=16, width=Inches(7.0))
    add_body(slide, Inches(1.7), Inches(4.0), [
        "SpectrumFM: a transformer foundation model for DESI spectra with an auxiliary redshift head "
        "(Huang DOE Genesis proposal).",
        "Ignition reproduced: peak val redshift-acc 14.86%; AR/TF ratio 0.73 ≈ NERSC reference 0.74.",
        "The load-bearing lever: the 4-way data mix (sv3+main × bright+dark) — not hyperparameters, scale, or seed.",
    ], size=16, width=Inches(7.0))
    add_picture_fit(slide, F_RS_DIV, Inches(7.7), Inches(1.5), Inches(5.2), Inches(4.4),
                    caption="Validation-loss trajectory — mix run (red) ignites first")
    add_footer(slide, "Full detail in reproductions/redshifty/papers/spectrumfm_status_2026-05-29.pptx")


def slide_integrity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "How these were done — and where to be careful")
    add_subtitle(slide, "Method", top=Inches(1.05), size=17, color=NAVY)
    add_body(slide, Inches(1.5), Inches(2.3), [
        "4 of 5 reproduced from the papers' methods sections — the Huang lens-finder code is not public.",
        "Public + collaboration data: DESI Legacy DR7/DR8, DESI DR1, HST/MAST.",
        "Local rig: 8× A16 + 2× L4 GPUs; every effort = numbered scripts + a LaTeX tech report (reproductions/*/papers/main.pdf).",
    ], size=16)
    add_subtitle(slide, "Honest caveats", top=Inches(4.05), size=17, color=ACCENT)
    add_body(slide, Inches(4.5), Inches(2.4), [
        "NeuraLens-catalog training leakage inflates absolute finder AUCs — so the faithful targets are the "
        "RELATIVE comparisons (shielded vs L18, DR9 vs DR7), which hold.",
        "Foundry I γ_EPL slope discrepancy (57% steeper) is open.",
        "Hsu 2025 visual-inspection grading is not reproduced (a team step).",
    ], size=15)
    add_footer(slide, "Stating the limits is what makes the matches credible.")


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What's next")
    add_body(slide, Inches(1.3), Inches(5.0), [
        "Finish Huang 2021 Phase 4c: merge DR8 shards → recovery by grade/threshold + leakage-aware breakdown + inspection viewer.",
        "Extend the lens-finder lineage: Storfer 2024 (DR9), Inchausti 2025 (DR10 dual-backbone + meta-learner ensemble).",
        "More GIGA-Lens application fits (Cikota 2023, Sheu 2024b) reusing the Foundry I harness.",
        "SpectrumFM: 20k-step extension (trajectory still descending at the 10k cap) + codecs Mamba+RFSQ tokenizer swap.",
    ], size=18)
    add_subtitle(slide, "None of these depend on NERSC compute to begin.",
                 top=Inches(6.3), size=16, color=ACCENT)


def slide_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Takeaways")
    add_body(slide, Inches(1.3), Inches(5.0), [
        "Five Huang-group results reproduced in ~a week — spanning Discover, Characterize, and the SpectrumFM foundation-model track.",
        "Method-from-paper reproduction works where code is closed: AUCs ≈0.99, θ_E to 3%, pair counts within 2%, ignition AR/TF 0.73 ≈ 0.74.",
        "The full local GPU + JAX + PyTorch environment is validated end-to-end.",
        "Ready to contribute to ongoing efforts: DR8 deployment, Foundry follow-ups, SpectrumFM Phase-I.",
    ], size=18)
    add_subtitle(slide, "From onboarding to production-grade reproductions in one sprint.",
                 top=Inches(6.3), size=16, color=NAVY)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Acknowledgments & questions")
    add_body(slide, Inches(1.3), Inches(4.0), [
        "Xiaosheng Huang (PI, USF) — program direction; GIGA-Lens & spectroscopy review.",
        "A. Storfer, Y.-M. Hsu, J. Samuel — candidate catalogs, the spherimatch tool, the redshifty repo.",
        "DESI Collaboration & public archives — Legacy Surveys, DESI DR1, HST/MAST.",
    ], size=18)
    add_subtitle(slide, "Questions / discussion / collaboration welcome.",
                 top=Inches(6.0), size=20, color=NAVY)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,
        slide_program,
        slide_scoreboard,
        slide_h20_method,
        slide_h20_results,
        slide_h21_method,
        slide_h21_results,
        slide_fi_method,
        slide_fi_results,
        slide_hsu_method,
        slide_hsu_results,
        slide_spectrumfm,
        slide_integrity,
        slide_roadmap,
        slide_takeaways,
        slide_thanks,
    ]
    for fn in builders:
        fn(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}  ({OUTPUT.stat().st_size} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
