#!/usr/bin/env python3
"""
Build the SpectrumFM-status PowerPoint deck for the group meeting.

Wraps python-pptx to construct an 11-slide deck reporting the Phase-13
local reproduction of the redshifty Approach-A ignition, anchored to the
DOE Genesis Phase I milestones (verbatim quotes from the proposal).

Run:
    python tools/spectrumfm/build_deck.py
Output:
    reproductions/redshifty/papers/spectrumfm_status_2026-05-29.pptx
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_ROOT = Path(__file__).resolve().parents[2]
OUTPUT = REPO_ROOT / "reproductions" / "redshifty" / "papers" / "spectrumfm_status_2026-06-02.pptx"
DIVERSITY_FIG = REPO_ROOT / "experiments" / "runs" / "_comparisons" / "diversity_journey.png"
COMPARISON_FIG = REPO_ROOT / "experiments" / "runs" / "_comparisons" / "v1_v2_codecs_l4x2_comparison.png"


# 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
LEFT_MARGIN = Inches(0.5)
TOP_MARGIN = Inches(0.4)
CONTENT_W = SLIDE_W - 2 * LEFT_MARGIN

NAVY = RGBColor(0x18, 0x2C, 0x5B)
ACCENT = RGBColor(0xC6, 0x3A, 0x3A)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)


def set_text(textframe, text, size=24, bold=False, color=INK, align=None):
    textframe.clear()
    p = textframe.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, text, size=32):
    box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, CONTENT_W, Inches(0.8))
    set_text(box.text_frame, text, size=size, bold=True, color=NAVY)
    return box


def add_subtitle(slide, text, top=Inches(1.1), size=18, color=MUTED):
    box = slide.shapes.add_textbox(LEFT_MARGIN, top, CONTENT_W, Inches(0.5))
    set_text(box.text_frame, text, size=size, color=color)
    return box


def add_body(slide, top, height, bullets, size=18, indent_size=14):
    """bullets: list[str] or list[tuple(level, str)]."""
    box = slide.shapes.add_textbox(LEFT_MARGIN, top, CONTENT_W, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = ("• " if level == 0 else "– ") + text
        run.font.size = Pt(size if level == 0 else indent_size)
        run.font.color.rgb = INK
    return box


def add_quote(slide, top, height, text, size=18):
    box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.5), top,
                                   CONTENT_W - Inches(1.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "“" + text + "”"
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.color.rgb = NAVY
    return box


def add_footer(slide, text):
    box = slide.shapes.add_textbox(LEFT_MARGIN, SLIDE_H - Inches(0.5),
                                   CONTENT_W, Inches(0.3))
    set_text(box.text_frame, text, size=11, color=MUTED)


def add_table(slide, top, height, rows, col_widths_in, header_row=True, size=14):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        LEFT_MARGIN, top,
        Inches(sum(col_widths_in)), height,
    )
    table = table_shape.table
    for col_idx, w in enumerate(col_widths_in):
        table.columns[col_idx].width = Inches(w)
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(size)
            if r_idx == 0 and header_row:
                run.font.bold = True
                run.font.color.rgb = NAVY
            else:
                run.font.color.rgb = INK
    return table_shape


def slide_1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.6), CONTENT_W, Inches(1.5))
    set_text(box.text_frame,
             "Reproducing SpectrumFM Phase-10 Ignition\non Commodity GPUs",
             size=40, bold=True, color=NAVY)
    box2 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.4), CONTENT_W, Inches(0.6))
    set_text(box2.text_frame, "Phase 13 Internal Reproduction Report",
             size=22, color=MUTED)
    box3 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(5.3), CONTENT_W, Inches(0.4))
    set_text(box3.text_frame, "Greg Benson, USF  ·  29 May 2026  ·  SpectrumFM team",
             size=16, color=INK)


def slide_2_one_model_six_classes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Context — “one model, six classes”")
    add_quote(slide, Inches(1.2), Inches(1.6),
              "The central claim is 'one model, six classes': a single SpectrumFM "
              "with one set of weights will handle six DESI target classes at "
              "production-grade quality. — DOE Genesis Project Narrative, v7")
    add_table(slide, Inches(3.2), Inches(2.7), [
        ["Regime", "Classes", "Phase I approach"],
        ["Label-rich (Redrock z)", "LRG, ELG, QSO (extragalactic); MWS (stars)",
         "Pre-training + auxiliary redshift head"],
        ["Extensibility", "LBGs, LAEs (extragalactic)",
         "Few-shot fine-tuning, ≤ 5,000 examples per class"],
    ], col_widths_in=[2.3, 4.7, 5.0])
    add_subtitle(slide, "This work validates the Phase-I pre-training + auxiliary "
                        "redshift head architecture at small scale on commodity GPUs.",
                 top=Inches(6.2), size=14, color=MUTED)
    add_footer(slide, "DOE Genesis SpectrumFM proposal v7 (Huang PI / USF, LBNL, NOIRLab, FSU)")


def slide_3_milestones(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase I Milestones — where this work lands")
    add_table(slide, Inches(1.1), Inches(3.0), [
        ["Timeline", "Focus area", "Measurable milestone"],
        ["Months 1–3", "Data Curation and Architecture",
         "DESI corpus and VI labels assembled. Architecture (auxiliary redshift "
         "head, physical prior regularization, MSM) implemented and validated on "
         "small-scale runs."],
        ["Months 4–6", "Pre-training from Scratch and Label-Rich Performance",
         "Full-scale pre-training on Perlmutter executed; scaling curves "
         "established. Objective 2a complete."],
        ["Months 7–9", "Extensibility, Alignment, and Validation",
         "Human Alignment Training executed. Objectives 1, 2b, and 3 complete. "
         "Phase I report and Phase II proposal submitted."],
    ], col_widths_in=[1.8, 2.7, 7.5], size=12)
    add_subtitle(slide,
                 "Months 1–3 explicitly call for “small-scale” validation of the "
                 "auxiliary redshift head — this is exactly that.",
                 top=Inches(5.3), size=16, color=ACCENT)
    add_subtitle(slide,
                 "Also de-risks the small-scale-runs requirement before ERCAP allocation lands.",
                 top=Inches(5.9), size=14, color=MUTED)
    add_footer(slide, "Quotes verbatim from proposals/doe_genesis_spectrumfm_project_narrative_v7.docx, Section 4")


def slide_4_tooling(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Local toolchain — Track 2 harness")
    add_subtitle(slide, "Yaml-driven experiment runner, parses both redshifty and codecs metric streams")
    add_body(slide, Inches(1.7), Inches(4.5), [
        "tools/spectrumfm/exp_run.py — yaml spec → subprocess + tee + metric parse → metrics.jsonl + summary.md",
        "tools/spectrumfm/exp_analyze.py — --digest <run> and --compare <runs...> (markdown table + PNG plot)",
        "tools/spectrumfm/download_desi_subset.py — public DESI mirror, --release {edr,dr1}, resumable",
        (1, "Drives both repos via one interface, mirroring redshifty/RESEARCH_LOG.md format"),
        "Five Approach-A arms (broken, phase10, xlarge, seed sweep, mix) ran through this harness unchanged.",
        (1, "Outputs preserved under experiments/runs/<run_id>/ with spec.yaml + stdout.log + metrics.jsonl + summary.md"),
    ], size=16)
    add_footer(slide, "~350 lines of Python total for the harness.")


def slide_5_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Data — local DESI mix")
    add_table(slide, Inches(1.1), Inches(2.4), [
        ["Survey × program", "Production", "Pixels", "Disk"],
        ["sv3 × bright", "EDR/fuji",  "373",  "304 GiB"],
        ["sv3 × dark",   "EDR/fuji",  "375",  "356 GiB"],
        ["main × bright","DR1/iron",  "192",  " 36 GiB"],
        ["main × dark",  "DR1/iron",  "197",  " 60 GiB"],
        ["TOTAL",        "",          "1,137","757 GiB"],
    ], col_widths_in=[3.5, 2.5, 2.5, 3.5], size=14)
    add_subtitle(slide, "1,817,790 raw spectra; 480k after quality cut "
                        "(ZCAT_PRIMARY ∧ COADD_FIBERSTATUS=0 ∧ OBJTYPE='TGT').",
                 top=Inches(4.3), size=14)
    add_subtitle(slide,
                 "NERSC reference: 200 healpix / 394k spectra across the SAME 4 combinations.",
                 top=Inches(5.0), size=14, color=ACCENT)
    add_subtitle(slide,
                 "Our local set has 5.7× more pixels and 4.6× more raw spectra than the reference.",
                 top=Inches(5.6), size=14, color=MUTED)
    add_footer(slide, "Pulled by tools/spectrumfm/download_desi_subset.py over ~3 h.")


def slide_6_tokenizer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "V1 tokenizer — local reproduction")
    add_body(slide, Inches(1.2), Inches(2.8), [
        "Architecture: ConvNeXt-V2 encoder/decoder + Look-Up-Free Quantizer (LFQ, dim=10, codebook=1024).",
        "Trained: 15,000 steps, batch 16, lr 3e-4, AMP enabled, one NVIDIA L4.",
        "Result: val_total = 1.70, val_recon = 1.38 at step 13,000 (best ckpt).",
        "NERSC reference: val_recon = 1.35 at step 16,500 on 200 healpix / 394k spectra.",
    ], size=18)
    add_subtitle(slide,
                 "Matched the NERSC baseline within 0.03 on val_recon — on roughly half the training corpus.",
                 top=Inches(5.0), size=18, color=ACCENT)
    add_footer(slide, "Frozen as the spectrum tokenizer for every Approach-A run that follows.")


def slide_7_diagnostic(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Diagnostic journey — 4 hypotheses for the ignition gap")
    add_table(slide, Inches(1.0), Inches(4.0), [
        ["Hypothesis", "Status", "Evidence"],
        ["Hparam mismatch (Phase 9 batch/lr × Phase 10 mask)",
         "Necessary but not sufficient",
         "Fix doubled peak val_z_acc 3.9 → 8.2%"],
        ["Data scale (need ≥394k spectra)",
         "Ruled out",
         "729k matched 219k exactly (6–8% plateau)"],
        ["Random-seed luck",
         "Ruled out",
         "4-seed sweep: 3.8%–8.8% range (1–3pp std)"],
        ["Data MIX (sv3+main × bright+dark)",
         "CONFIRMED",
         "First arm to ignite, sustained ≥10%, AR/TF 0.73"],
    ], col_widths_in=[4.8, 2.8, 4.6], size=12)
    add_subtitle(slide,
                 "The buried clue: published log line “200 healpix files (sv3+main, bright+dark)” "
                 "— easy to misread as 200 sv3-bright pixels.",
                 top=Inches(5.4), size=14, color=MUTED)
    add_footer(slide, "Each row took 5–8 h of training to confirm or refute.")


def slide_8_ignition(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Ignition achieved — mix-run val trajectory")
    # smaller table on the left
    add_table(slide, Inches(0.5), Inches(4.5), [
        ["step", "val_z_acc", "val_loss_redshift"],
        ["500",  "0.028", "4.88"],
        ["3000", "0.015", "4.50"],
        ["5000", "0.034", "4.24"],
        ["6500", "0.055", "4.08"],
        ["7500", "0.081", "3.97"],
        ["8500", "0.092", "3.81"],
        ["9000", "0.108", "3.73"],
        ["9500", "0.149", "3.69"],
    ], col_widths_in=[1.0, 1.5, 1.8], size=12)
    # figure on the right
    slide.shapes.add_picture(str(DIVERSITY_FIG),
                             Inches(5.5), Inches(1.3),
                             width=Inches(7.5))
    add_subtitle(slide,
                 "Mix run (red): first to cross sustained 10%; reaches lowest val_loss (190.7).",
                 top=Inches(6.4), size=14, color=ACCENT)


def slide_9_pass(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Pass criteria from the plan — all three met")
    add_table(slide, Inches(1.0), Inches(3.5), [
        ["Criterion", "Target", "Achieved", "Verdict"],
        ["val_redshift_acc sustained ≥ 10%", "≥10% across 2+ vals",
         "10.85% @ step 9000, 14.86% @ step 9500", "PASS"],
        ["val_loss_redshift cumulative drop", "≥ 1.0",
         "4.88 → 3.69 (drop 1.19)", "PASS"],
        ["AR / TF ratio at peak", "≥ 0.5", "7.96% / 10.85% = 0.73", "PASS"],
    ], col_widths_in=[3.5, 2.5, 3.5, 1.5], size=12)
    add_subtitle(slide,
                 "Best val_loss = 190.67 (vs all prior runs' 200–225 floor).",
                 top=Inches(4.8), size=16)
    add_subtitle(slide,
                 "AR/TF ratio 0.73 vs NERSC reference 0.74 — quantitatively matches the published result.",
                 top=Inches(5.5), size=16, color=ACCENT)


def slide_phase14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 14–15 — tokenizer comparison (V1 / V2 / codecs)")
    add_table(slide, Inches(0.5), Inches(1.95), [
        ["arm", "tokenizer", "TF z_acc", "AR z_acc", "codebook", "fair good-z"],
        ["V1", "ConvNeXt+LFQ (256-lvl z)", "0.55", "0.57", "5.00 bits", "0.192"],
        ["V2 +skips", "+U-Net skips (1024)", "0.00", "0.00", "0.00 bits (collapsed)", "—"],
        ["V2 no-skip", "+tophat/entropy (1024)", "0.50", "0.48", "5.24 bits", "0.192"],
        ["codecs", "Mamba3+RFSQ (256-lvl z)", "0.53", "0.52", "6.27 bits", "—"],
    ], col_widths_in=[1.3, 3.3, 1.2, 1.2, 2.3, 1.4], size=12)
    slide.shapes.add_picture(str(COMPARISON_FIG), Inches(2.4), Inches(3.0), width=Inches(8.5))
    add_subtitle(slide,
                 "Three structurally different healthy tokenizers (ConvNeXt+LFQ, +tophat, "
                 "Mamba3+RFSQ) all converge to ~0.50–0.55. Tokenizer architecture is NOT the "
                 "bottleneck once the codebook is healthy; V2+skips collapses (codes bypassed).",
                 top=Inches(6.55), size=12, color=ACCENT)


def slide_phase16_perclass(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 16 — per-class redshift parity vs Redrock (Metric 1)")
    add_subtitle(slide, "Honest encoder-masked |Δz|/(1+z); ZWARN=0 held-out set; "
                 "pass = catastrophic rate within 5pp of zero", top=Inches(1.05), size=15)
    add_table(slide, Inches(0.5), Inches(3.1), [
        ["class", "N", "catastrophic rate", "median |Δz|/(1+z)", "verdict"],
        ["MWS (stars)", "289", "2.1%", "0.0002", "PASS"],
        ["BGS", "346", "95.7%", "0.0361", "FAIL"],
        ["LRG", "202", "97.5%", "0.0630", "FAIL"],
        ["ELG", "366", "97.8%", "0.1106", "FAIL"],
        ["QSO *", "74", "97.3%", "0.1782", "FAIL"],
    ], col_widths_in=[2.6, 1.2, 3.2, 3.2, 1.8], size=14)
    add_subtitle(slide,
                 "Only stars reach good-z parity. Failure ordering MWS≪BGS<LRG<ELG<QSO "
                 "tracks spectroscopic difficulty exactly → the metric is real. Coarse z is "
                 "learned (60% within |Δz|/(1+z)<0.05); DESI precision is not. V1≈V2 no-skip.",
                 top=Inches(5.6), size=13, color=ACCENT)
    add_footer(slide, "tools/spectrumfm/eval_per_class.py (+ desi_targets.py); decode vs SPECTYPE 98.7% agree. * QSO Redrock-only comparator.")


def slide_phase16_probe(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 16 — frozen-encoder six-class probe")
    add_subtitle(slide, "Linear probe on frozen mean-pooled encoder features "
                 "(redshift token masked); healpix-disjoint split", top=Inches(1.05), size=15)
    add_table(slide, Inches(0.5), Inches(2.6), [
        ["metric", "V1", "V2 no-skip", "no-skill baseline"],
        ["macro-F1", "0.813", "0.847", "0.097"],
        ["accuracy", "0.858", "0.886", "0.321"],
        ["per-class F1 (V2)", "MWS .95 / BGS .92 / ELG .90", "LRG .87 / QSO .61", "—"],
    ], col_widths_in=[2.6, 3.6, 3.0, 2.6], size=13)
    add_subtitle(slide,
                 "The encoder learns WHAT an object is (macro-F1 0.81–0.85, +0.72 over no-skill) "
                 "even where it cannot pin HOW FAR (Metric 1). Classification capability and "
                 "redshift precision are DISTINCT — the precision gap is not a representation problem.",
                 top=Inches(5.0), size=14, color=ACCENT)
    add_subtitle(slide, "Only confusions are physical: QSO↔ELG (emission-line), LRG↔BGS (red galaxies).",
                 top=Inches(6.1), size=13, color=MUTED)
    add_footer(slide, "tools/spectrumfm/probe_six_class.py; supports the proposal's “one encoder, six classes” claim.")


def slide_phase16_eqprior(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 16 — physical-prior ablation: redshift equivariance")
    add_subtitle(slide, "Self-supervised: shift a spectrum by δz → require predicted z to move by δz "
                 "(flag-gated, default-off)", top=Inches(1.05), size=15)
    add_table(slide, Inches(0.5), Inches(2.7), [
        ["equivariance response  d E[z]/d(δz)  (ideal 1.0)", "V1 control", "eqprior"],
        ["δz = +0.05", "0.00", "0.35"],
        ["δz = +0.10", "0.10", "0.27"],
        ["δz = +0.20", "0.15", "0.14"],
    ], col_widths_in=[6.0, 2.6, 2.6], size=14)
    add_subtitle(slide,
                 "Prior WORKS mechanically (non-equivariant 0.00 → 0.35 at small shifts) but per-class "
                 "catastrophic rate is UNCHANGED (aggregate good-z 24.1% → 24.5%).",
                 top=Inches(5.0), size=14, color=ACCENT)
    add_subtitle(slide,
                 "Principled null: equivariance is shift-CONSISTENCY; galaxy |Δz| (~0.04–0.18) is off by "
                 "10–50× the 0.0033 threshold. Consistency cannot fix absolute error → SCALE is the lever.",
                 top=Inches(6.0), size=13, color=NAVY)
    add_footer(slide, "nersc/physical_priors.py; 15k-step matched treatment vs V1 control, both L4s.")


def slide_10_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What's next")
    add_body(slide, Inches(1.2), Inches(5.0), [
        "DONE: 2×L4 rerun (correct PCI_BUS_ID mapping, bf16, eff batch 64) → V1 55% TF / 57% AR",
        (1, "3.7× the A16 mix run; honest AR matches the NERSC reference (55%)"),
        "DONE: V2 tokenizer revisited — skip-bypass codebook collapse diagnosed and fixed",
        (1, "Skip-free V2 ties V1: codebook entropy, not reconstruction loss, predicts usefulness"),
        "DONE: codecs (Mamba3+RFSQ) tokenizer swap — ties V1 (0.53 TF / 0.52 AR)",
        (1, "Tokenizer architecture is not the bottleneck; codebook health gates usability"),
        "DONE: Phase-16 evaluation framework — per-class Metric-1 eval + frozen six-class probe",
        (1, "Metric 1 not met at scale (only stars); encoder learns class strongly (macro-F1 0.81–0.85)"),
        "DONE: physical-prior ablation (redshift equivariance) — works mechanically, no precision gain",
        (1, "Equivariance & absolute precision are decoupled → scale, not aux losses, is the lever"),
        "Next: local data×model-size scaling ladder (per-point metric = Metric-1 cat. rate + probe-F1)",
        (1, "Then the full-DR1 / 250M–1B-param NERSC spec, with eval_per_class.py as the completion gate"),
        "Open PRs: cosmologyfoundation/redshifty#1 (+ multi-GPU fixes), codecs#1",
    ], size=14)
    add_footer(slide, "Local prototyping done; precision now requires NERSC-scale compute (the proposal's Phase-I premise).")


def slide_11_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Acknowledgments")
    add_body(slide, Inches(1.2), Inches(4.5), [
        "Jonathan Samuel (USF) — author of the redshifty Approach-A repo and the PRODUCTION_RUN_PLAN.md / RESEARCH_LOG.md that this reproduction works against.",
        "Xiaosheng Huang (USF) — SpectrumFM PI; spectroscopy-side architecture review.",
        "Stephen Bailey (LBNL) — Redrock + DESI corpus expertise that grounds the 4-way data-mix interpretation.",
        "DESI Collaboration — public EDR/DR1 data tier.",
        "Anthropic / Claude Code — the diagnostic-ladder + harness implementation pair-programmer.",
    ], size=16)
    add_subtitle(slide, "Questions / discussion / PR review welcome.",
                 top=Inches(6.4), size=18, color=NAVY)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs)
    slide_2_one_model_six_classes(prs)
    slide_3_milestones(prs)
    slide_4_tooling(prs)
    slide_5_data(prs)
    slide_6_tokenizer(prs)
    slide_7_diagnostic(prs)
    slide_8_ignition(prs)
    slide_9_pass(prs)
    slide_phase14(prs)
    slide_phase16_perclass(prs)
    slide_phase16_probe(prs)
    slide_phase16_eqprior(prs)
    slide_10_next(prs)
    slide_11_thanks(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}  ({OUTPUT.stat().st_size} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
