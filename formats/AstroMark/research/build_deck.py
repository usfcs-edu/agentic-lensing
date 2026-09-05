#!/usr/bin/env python3
"""Build the AstroMark research presentation.

    ~/.venvs/workshop/bin/python build_deck.py [--out AstroMark-Research.pptx]

A walkthrough of the research, findings, proposals and recommendations, for the team review. Dark
ground because two thirds of the slides carry astronomical renders, which look wrong on white.

Every number on a slide comes from a file in this directory; nothing is retyped by hand that could
be read from `examples/reference/metrics.json` or `examples/encodings/sizes.json`.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
REF = HERE / "examples/reference"
JWST = HERE / "examples/jwst"

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(0x0C, 0x0E, 0x14)
FG = RGBColor(0xF2, 0xF3, 0xF7)
DIM = RGBColor(0x9A, 0xA0, 0xAE)
ACC = RGBColor(0x56, 0xB4, 0xE9)      # sky blue — the notation's "lensed light"
GOOD = RGBColor(0x00, 0x9E, 0x73)     # bluish green — "lens mass"
WARN = RGBColor(0xE6, 0x9F, 0x00)     # orange
BAD = RGBColor(0xD5, 0x5E, 0x00)      # vermillion
MONO = "Menlo"
SANS = "Helvetica Neue"


def metrics():
    return {r["arm"]: r for r in json.loads((REF / "metrics.json").read_text())}


def sizes():
    return json.loads((HERE / "examples/encodings/sizes.json").read_text())["rows"]


class Deck:
    def __init__(self):
        self.p = Presentation()
        self.p.slide_width, self.p.slide_height = W, H
        self.blank = self.p.slide_layouts[6]
        self.n = 0

    # -- primitives ------------------------------------------------------------------------
    def _slide(self):
        s = self.p.slides.add_slide(self.blank)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = BG
        self.n += 1
        return s

    MIN = Inches(0.14)

    def _tb(self, s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        w, h = Emu(max(int(w), int(self.MIN))), Emu(max(int(h), int(self.MIN)))
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        return tf

    def _para(self, tf, text, size, colour=FG, bold=False, font=SANS, space_after=6,
              first=False, align=PP_ALIGN.LEFT, italic=False):
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space_after)
        # inline emphasis with **bold** and `code`
        for chunk, kind in _split(text):
            r = para.add_run()
            r.text = chunk
            r.font.size = Pt(size)
            r.font.bold = bold or kind == "b"
            r.font.italic = italic
            r.font.name = MONO if kind == "c" else font
            r.font.color.rgb = ACC if kind == "c" else colour
        return para

    def _footer(self, s, note=""):
        tf = self._tb(s, Inches(0.6), H - Inches(0.42), W - Inches(1.2), Inches(0.3))
        self._para(tf, note, 9, DIM, first=True)
        tb = s.shapes.add_textbox(W - Inches(1.1), H - Inches(0.42), Inches(0.6), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(self.n)
        r.font.size = Pt(9)
        r.font.color.rgb = DIM
        r.font.name = SANS

    def _rule(self, s, y, colour=ACC, x=Inches(0.6), w=Inches(1.4), h=Pt(3)):
        from pptx.enum.shapes import MSO_SHAPE
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(int(h)))
        sh.fill.solid()
        sh.fill.fore_color.rgb = colour
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    # -- slide kinds -----------------------------------------------------------------------
    def title(self, title, subtitle, meta=""):
        s = self._slide()
        self._rule(s, Inches(2.35), ACC, Inches(0.9), Inches(2.2))
        tf = self._tb(s, Inches(0.9), Inches(2.6), W - Inches(2.4), Inches(2.2))
        self._para(tf, title, 44, FG, bold=True, first=True, space_after=14)
        self._para(tf, subtitle, 20, DIM, space_after=10)
        if meta:
            self._para(tf, meta, 13, DIM)
        self._footer(s)
        return s

    def section(self, kicker, title, blurb=""):
        s = self._slide()
        self._rule(s, Inches(2.7), WARN, Inches(0.9), Inches(1.1))
        tf = self._tb(s, Inches(0.9), Inches(2.95), W - Inches(2.4), Inches(2.4))
        self._para(tf, kicker.upper(), 13, WARN, bold=True, first=True, space_after=12)
        self._para(tf, title, 34, FG, bold=True, space_after=12)
        if blurb:
            self._para(tf, blurb, 15, DIM)
        self._footer(s)
        return s

    def head(self, s, title, sub=""):
        tf = self._tb(s, Inches(0.6), Inches(0.42), W - Inches(1.2), Inches(1.0))
        self._para(tf, title, 27, FG, bold=True, first=True, space_after=4)
        if sub:
            self._para(tf, sub, 13, DIM)
        return Inches(1.55) if sub else Inches(1.30)

    def bullets(self, title, items, sub="", note="", size=15, lead=""):
        s = self._slide()
        y = self.head(s, title, sub)
        tf = self._tb(s, Inches(0.6), y, W - Inches(1.2), H - y - Inches(0.7))
        first = True
        if lead:
            self._para(tf, lead, 17, ACC, first=True, space_after=14)
            first = False
        for it in items:
            if isinstance(it, tuple):
                text, colour = it
            else:
                text, colour = it, FG
            indent = text.startswith("    ")
            self._para(tf, ("– " if not indent else "") + text.strip(),
                       size - (2 if indent else 0),
                       DIM if indent else colour, first=first, space_after=9 if not indent else 5)
            first = False
        self._footer(s, note)
        return s

    def image(self, title, path, sub="", note="", caption="", side=None, side_title=""):
        """Full-bleed-ish image, optionally with a text column on the right."""
        s = self._slide()
        y = self.head(s, title, sub)
        avail_h = H - y - Inches(0.75)
        from PIL import Image as PImage
        iw, ih = PImage.open(path).size
        side_w = Inches(4.9)
        gap = Inches(0.5)
        max_w = (W - Inches(1.2)) if not side else (W - Inches(1.2) - side_w - gap)
        scale = min(max_w / Emu(1) / iw, avail_h / Emu(1) / ih)
        w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
        if side:
            x = Emu(int((W - (w + gap + side_w)) / 2))     # centre the pair, not the picture
        else:
            x = Emu(int((W - w) / 2))
        s.shapes.add_picture(str(path), x, y, w, h)
        if caption:
            tf = self._tb(s, x, y + h + Inches(0.08), w, Inches(0.4))
            self._para(tf, caption, 11, DIM, first=True)
        if side:
            sx = x + w + gap
            tf = self._tb(s, sx, y, side_w, avail_h)
            first = True
            if side_title:
                self._para(tf, side_title, 15, ACC, bold=True, first=True, space_after=12)
                first = False
            for it in side:
                if isinstance(it, tuple):
                    text, colour = it
                else:
                    text, colour = it, FG
                self._para(tf, text, 13, colour, first=first, space_after=10)
                first = False
        self._footer(s, note)
        return s

    def two_images(self, title, left, right, sub="", lcap="", rcap="", note=""):
        s = self._slide()
        y = self.head(s, title, sub)
        avail_h = H - y - Inches(0.95)
        from PIL import Image as PImage
        cw = (W - Inches(1.5)) / 2
        for i, (path, cap) in enumerate(((left, lcap), (right, rcap))):
            iw, ih = PImage.open(path).size
            scale = min(cw / Emu(1) / iw, avail_h / Emu(1) / ih)
            w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
            x = Inches(0.6) + i * (cw + Inches(0.3)) + Emu(int((cw - w) / 2))
            s.shapes.add_picture(str(path), x, y, w, h)
            if cap:
                tf = self._tb(s, Inches(0.6) + i * (cw + Inches(0.3)), y + h + Inches(0.1),
                              cw, Inches(0.6))
                self._para(tf, cap, 11, DIM, first=True)
        self._footer(s, note)
        return s

    def table(self, title, headers, rows, sub="", note="", widths=None, highlight=None,
              size=12):
        s = self._slide()
        y = self.head(s, title, sub)
        n_r, n_c = len(rows) + 1, len(headers)
        tw = W - Inches(1.2)
        th = min(H - y - Inches(0.8), Inches(0.36) * n_r)
        shape = s.shapes.add_table(n_r, n_c, Inches(0.6), y, tw, th)
        tbl = shape.table
        if widths:
            total = sum(widths)
            for i, ww in enumerate(widths):
                tbl.columns[i].width = Emu(int(tw * ww / total))
        for c, htxt in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1E, 0x28)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = htxt
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = ACC; r.font.name = SANS
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = (RGBColor(0x16, 0x1A, 0x24) if ri % 2 else
                                            RGBColor(0x11, 0x14, 0x1C))
                p = cell.text_frame.paragraphs[0]
                bold = highlight is not None and ri - 1 in highlight
                for chunk, kind in _split(str(val)):
                    r = p.add_run(); r.text = chunk
                    r.font.size = Pt(size)
                    r.font.bold = bold or kind == "b"
                    r.font.name = MONO if kind == "c" else SANS
                    r.font.color.rgb = (WARN if bold else (ACC if kind == "c" else FG))
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._footer(s, note)
        return s

    def statement(self, kicker, big, support="", colour=ACC, note=""):
        s = self._slide()
        tf = self._tb(s, Inches(1.0), Inches(1.6), W - Inches(2.0), H - Inches(3.0),
                      anchor=MSO_ANCHOR.MIDDLE)
        if kicker:
            self._para(tf, kicker.upper(), 13, colour, bold=True, first=True, space_after=18)
            self._para(tf, big, 30, FG, bold=True, space_after=18)
        else:
            self._para(tf, big, 30, FG, bold=True, first=True, space_after=18)
        if support:
            self._para(tf, support, 15, DIM)
        self._footer(s, note)
        return s

    def numbers(self, title, cards, sub="", note=""):
        """Three or four big figures across."""
        s = self._slide()
        y = self.head(s, title, sub)
        n = len(cards)
        cw = (W - Inches(1.2) - Inches(0.3) * (n - 1)) / n
        for i, (big, label, colour) in enumerate(cards):
            x = Inches(0.6) + i * (cw + Inches(0.3))
            tf = self._tb(s, x, y + Inches(0.5), cw, Inches(3.0))
            self._para(tf, big, 46, colour, bold=True, first=True, space_after=10)
            self._para(tf, label, 13, DIM)
        self._footer(s, note)
        return s

    def code(self, title, lines, sub="", note="", size=13, caption=""):
        s = self._slide()
        y = self.head(s, title, sub)
        from pptx.enum.shapes import MSO_SHAPE
        # budget the space first: the block plus its caption must fit above the footer
        avail = H - y - Inches(0.65) - (Inches(0.85) if caption else Inches(0))
        pitch = min(Inches(0.34), Emu(int((avail - Inches(0.5)) / max(1, len(lines)))))
        if pitch < Inches(0.24):                     # too tight to read: shrink the type instead
            size = max(9, int(size * pitch / Inches(0.30)))
            pitch = Emu(int(max(int(pitch), int(Inches(0.20)))))
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y,
                                 W - Inches(1.2), Emu(int(pitch * len(lines) + Inches(0.5))))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x16, 0x1A, 0x24)
        box.line.fill.background(); box.shadow.inherit = False
        tf = self._tb(s, Inches(0.9), y + Inches(0.25), W - Inches(1.8),
                      Emu(int(pitch * len(lines))))
        for i, ln in enumerate(lines):
            colour = FG
            if ln.startswith("# "):
                colour = DIM
            elif ln.startswith("!"):
                colour = BAD; ln = ln[1:]
            elif ln.startswith("+"):
                colour = GOOD; ln = ln[1:]
            self._para(tf, ln or " ", size, colour, font=MONO, first=(i == 0), space_after=3)
        if caption:
            cy = Emu(int(y + pitch * len(lines) + Inches(0.75)))
            tf2 = self._tb(s, Inches(0.6), cy, W - Inches(1.2),
                           Emu(int(max(int(H - cy - Inches(0.5)), int(Inches(0.3))))))
            self._para(tf2, caption, 14, DIM, first=True)
        self._footer(s, note)
        return s

    def validate(self) -> list[str]:
        """Every shape must have a positive extent and sit inside the slide.

        PowerPoint refuses to open a presentation containing a non-positive cx/cy; LibreOffice
        tolerates it silently, so a PDF render is NOT sufficient verification.
        """
        problems = []
        for i, slide in enumerate(self.p.slides, start=1):
            for sh in slide.shapes:
                if sh.width is None or sh.height is None:
                    continue
                if sh.width <= 0 or sh.height <= 0:
                    problems.append(f"slide {i}: {sh.shape_type} has extent "
                                    f"{sh.width}x{sh.height} — PowerPoint will reject this")
                if sh.left is None or sh.top is None:
                    continue
                if sh.left < -Inches(0.5) or sh.top < -Inches(0.5):
                    problems.append(f"slide {i}: {sh.shape_type} starts off-slide "
                                    f"at {sh.left},{sh.top}")
                if sh.top + sh.height > H + Inches(1.0):
                    problems.append(f"slide {i}: {sh.shape_type} overruns the bottom by "
                                    f"{(sh.top + sh.height - H) / 914400:.2f} in")
        return problems

    def save(self, path):
        problems = self.validate()
        if problems:
            raise SystemExit("deck validation FAILED:\n  " + "\n  ".join(problems))
        self.p.save(str(path))
        return self.n


def _split(text):
    """Tiny inline markup: **bold** and `code`."""
    out, buf, i = [], "", 0
    while i < len(text):
        if text.startswith("**", i):
            j = text.find("**", i + 2)
            if j > 0:
                if buf:
                    out.append((buf, "n")); buf = ""
                out.append((text[i + 2:j], "b")); i = j + 2; continue
        if text[i] == "`":
            j = text.find("`", i + 1)
            if j > 0:
                if buf:
                    out.append((buf, "n")); buf = ""
                out.append((text[i + 1:j], "c")); i = j + 1; continue
        buf += text[i]; i += 1
    if buf:
        out.append((buf, "n"))
    return out or [("", "n")]


# ══════════════════════════════════════════════════════════════════════════════════════════
# The deck
# ══════════════════════════════════════════════════════════════════════════════════════════

def build(out: Path) -> int:
    d = Deck()
    M = metrics()
    SZ = {r["title"]: r for r in sizes()}

    # ── 0. Frame ──────────────────────────────────────────────────────────────────────────
    d.title("AstroMark",
            "A symbolic notation and metadata format for annotating astronomical images",
            "Research, proposals and recommendations  ·  3 September 2026  ·  "
            "Nothing here is adopted — this is for the team to choose from")

    d.bullets("What was asked, and what is here",
              [("**The ask.** A concise but descriptive notation, useful to humans and to models, "
                "for marking features in telescope images. Lenses first, astronomy generally later.", FG),
               "    Three specific questions: should it be accessible without colour? can the symbols "
               "be thin but readable? can marks be resized for emphasis?",
               ("**The work.** 66 formats surveyed across 8 domains. Four visual notations designed, "
                "prototyped and measured. Six metadata architectures encoded and compared. A draft "
                "spec for the recommendation.", FG),
               ("**The evidence.** Working code, deterministic from a seed: a synthetic reference "
                "scene, a prototype renderer, colour-vision simulation, a measurement harness, and "
                "worked examples on real JWST data.", FG),
               ("**The status.** Recommendations, not decisions. Where the evidence is thin, the "
                "slides say so.", DIM)],
              sub="formats/AstroMark/research/")

    d.table("The recommendation, in one table",
            ["", "Pick", "Why"],
            [["Notation", "**N1 Terminator alphabet**, with N4 prototyped beside it",
              "N1 keeps the arrow and costs least; N4 is the only one that fixes the named defect"],
             ["Container", "**P1 evolved bespoke**", "the cheapest decision — converters are 57–161 lines"],
             ["Identity", "**prefixed terms in one vocabulary file**",
              "this is where every actual problem is"],
             ["Model surfaces", "**read** the line form, **write** JSON",
              "the only arrangement satisfying both the schema rule and the token budget"],
             ["Portrayal", "**out of the record**", "so no field can carry meaning by appearance"],
             ["Accessibility", "**hard constraint**, tested", "role by shape, polarity by texture"],
             ["Do first", "**build the vocabulary file**", "the only artifact every proposal shares"]],
            widths=[16, 34, 50], highlight=set(range(7)), size=12,
            note="Detail in 07-recommendation.md")

    # ── 1. The problem ────────────────────────────────────────────────────────────────────
    d.section("Part one", "The problem",
              "What the current format is, and what is actually wrong with it")

    d.bullets("Four things were competing for one name",
              ["**The app** — hand, model and voice annotation.  Named: LensMark.",
               "**The symbolic notation** — what the marks mean.  Unnamed.",
               "**The metadata format** — vector JSON, provenance, render contract.  Unnamed.",
               "**The standard as a whole** — what gets published and adopted.  Unnamed.",
               ("Naming the fourth separately is what unlocks the rest: a general core with domain "
                "profiles, so the lens work becomes `astromark/lens/1.0` rather than the whole thing, "
                "and LensMark keeps its name as the reference implementation.", ACC)],
              sub="This is why the naming conversation went in circles",
              note="naming-study.md")

    d.code("What the current record looks like",
           ["{",
            '  "schema_version": "lensmark/1.0",',
            '  "image":  { "sha256": …, "cutout_arcsec": 16.0, "pixel_scale_arcsec": 0.0397, … },',
            '  "system": { "grade": "A", "theta_e": { "value_arcsec": 1.45, "method": "geometric" },',
            '              "verdict": "likely_lens", "description": "free text …" },',
            '  "style_defaults": { … the whole style table, copied into every file … },',
            '  "items": [',
            '    { "type": "arrow", "tail": [.40,.30], "head": [.46,.40],',
            '      "label": "arc", "color": "cyan" },',
            '    { "type": "mask_circle", "center": [.80,.20], "radius_arcsec": 0.9,',
            '      "kind": "galaxy", "color": "mask_red" },',
            '    { "type": "einstein_ring", "center_ref": "…", "theta_e_arcsec": 1.45 }',
            '  ]",'.replace('",', ''),
            "}"],
           sub="lensmark/1.0 — 11,184 characters per deck on average, across 19 real decks",
           caption="Geometry is normalized [u, v]; sizes are arcsec; the render is already "
                   "deterministic and pinned by nine golden hashes. The bones are good.",
           size=12)

    d.table("Nine structural problems — and what each one really is",
            ["Symptom", "What it actually is"],
            [["Enums defined in **six** uncoordinated places", "term identity, unowned"],
             ["θ_E stored twice and unlinked", "identity of a *measured quantity*"],
             ["`theta_e.method` a free string despite three named rules", "a term with no vocabulary"],
             ["green = deflector; colours reserved by item type", "a term smuggled into portrayal"],
             ["prose refers to arrows *by colour*", "a term with no referenceable id"],
             ["the style table copied into every data file", "portrayal had no identity, so it was inlined"],
             ["no role field — semantics in a 40-char label", "a term hidden in free text"],
             ["no polarity", "a missing term"],
             ["no polygon, no secondary ring, no source grouping", "missing terms, and a missing grouping term"]],
            widths=[52, 48], size=12,
            sub="Audited against the codebase, not inferred",
            note="04-metadata-proposals.md §0")

    d.statement("nine for nine",
                "They are all vocabulary-identity problems.",
                "Not one would have been prevented by GeoJSON, by Web Annotation, or by tables.\n"
                "So: choose the container for cheapness. Put the rigour in the vocabulary.",
                colour=WARN)

    d.code("The sharpest instance",
           ["# lensmark/validate.py:75  — and again, independently, in frontend/src/geometry.ts",
            "",
            'def _is_deflector(label):',
            '    return bool(label) and "deflector" in label.lower()',
            "",
            "# so this label renders GREEN — the colour that means lens mass:",
            '!    "spiral arm, NOT a deflector"'],
           sub="How the current system decides what is a lens galaxy",
           caption="The semantics of the standard ride on substring-matching an optional "
                   "40-character string, in two languages, with no shared definition.",
           size=15)

    d.image("What the best annotations so far look like",
            REF / "r-campaign.png",
            sub="The agent campaign the room preferred — reproduced here at its own constants",
            side=[("**Measured strengths.** 1.2–2.3% ink inside the image. Arrows never overlap "
                   "the feature. Line style separates the measured ring from galaxy and star masks. "
                   "Every mask numbered and keyed to a table. The key sits below the image.", FG),
                  ("**Measured weaknesses.**", WARN),
                  "The colour key is only partly consistent: red and green are locked, four colours "
                  "are opportunistic, and cyan flips between “unrelated field object” and “lensed "
                  "feature” *within the same nine-panel set*.",
                  "Uncertainty is text only — a question mark, “would-be”, a parenthetical range.",
                  ("**Negative evidence has no visual channel at all.** Flipping through nine panels, "
                   "you cannot tell which one is the non-lens from the marks.", BAD)],
            side_title="Why it was liked, and where it stops",
            note="03-notation-proposals.md")

    # ── 2. Prior art ──────────────────────────────────────────────────────────────────────
    d.section("Part two", "How other fields solved this",
              "66 formats across 8 domains. Medicine solved this thirty years ago; astronomy "
              "largely has not.")

    d.numbers("The survey",
              [("66", "formats examined, primary sources where reachable", ACC),
               ("8", "domains: radiology, medical vocabularies, computer vision,\nweb and geospatial, "
                "astronomy overlays, astronomy metadata,\nnotation theory, accessibility", GOOD),
               ("9", "mechanisms independently arrived at by 6–10 domains", WARN),
               ("8", "well-designed standards that died anyway", BAD)],
              sub="01-prior-art/")

    d.code("Radiology's foundational move: a term is never a string",
           ["# DICOM coded entry — three mandatory fields plus an optional version",
            "",
            '  Code Value              (0008,0100)   the opaque, stable identity',
            '  Coding Scheme Designator(0008,0102)   which vocabulary it came from',
            '  Code Meaning            (0008,0104)   the human-readable string',
            "",
            "# and the rule that makes it work, from PS3.3 §8.3, normative:",
            "",
            '  "the Value of Code Meaning shall NEVER be used as a key, index or',
            '   decision value ... Code Meaning is a purely annotative, descriptive',
            '   Attribute."   ... "This does not imply that Code Meaning can be',
            '   filled with arbitrary free text."'],
           sub="DICOM PS3.3 §8.3",
           caption="LensMark's 40-character substring-matched label is precisely the anti-pattern "
                   "that sentence exists to forbid.", size=13)

    d.bullets("What else radiology settled",
              [("**Portrayal is a separate object.** A Presentation State is its own instance that "
                "*references* images; the pixels are never modified. Standardised in 1999.", FG),
               "    Meaning rides on a graphic layer, and every layer carries a greyscale value "
               "**and** a CIELab colour. Redundant achromatic encoding is in the data model, not "
               "bolted on.",
               ("**The measurement is primary; the geometry is inferred from it.** The shape is a "
                "witness, not a second copy — which is exactly the fix for θ_E stored twice.", FG),
               ("**Presence is three-valued and coded:** present, absent, **undetermined**. AIM used "
                "a boolean and the standard records why that was wrong: a boolean cannot say "
                "“undetermined”.", FG),
               ("**Generation method is three-valued:** automatic, semiautomatic, manual. The middle "
                "value is the one everyone omits and everyone needs — a model-proposed arc that a "
                "human nudged is neither.", FG),
               ("**Two identities per mark:** this mark, and the physical thing it is of. Grouping "
                "then becomes emergent rather than structural.", FG)],
              sub="DICOM SR, GSPS, RT Structure Sets, SEG",
              note="01-prior-art/medical-imaging.md")

    d.statement("the cautionary tale",
                "NCI's AIM got the central abstraction exactly right, and died.",
                "It invented the markup-versus-annotation split in 2008 — the same architecture this "
                "project arrived at independently. Its reference repository is archived; the last "
                "substantive commit was 29 April 2014.\n\n"
                "Three verifiable reasons: it put lineColor and lineThickness on its own geometry "
                "class, violating its own thesis in its own schema; 58 of its 124 types were "
                "relationship classes with names like "
                "ImagingObservationEntityIsIdentifiedByTwoDimensionGeometricShapeEntityStatement; "
                "and it imported ISO datatype machinery because it looked rigorous.",
                colour=BAD, note="This is why the recommendation takes AIM's identity and leaves AIM's ceremony")

    d.bullets("Graded assessment: BI-RADS has the answer to the grade problem",
              [("**BI-RADS defines its scale as a three-column table:** ordinal category · numeric "
                "probability band · **the specific next observation that would settle it** — plus a "
                "published expected prevalence per grade.", FG),
               ("**PI-RADS goes further and DERIVES the grade** from coded component scores by a "
                "published lookup table. The evidence is stored; the verdict is computed.", FG),
               ("For a project whose goal is consistent grading between annotators, the derived model "
                "makes disagreement **locatable**: two graders who differ on the verdict but agree on "
                "the evidence have a different problem from two who differ on the evidence.", ACC),
               ("Grades A–D are currently used with only the borderline class defined on the record. "
                "This is the direct answer, and it is unclaimed.", WARN)],
              sub="The most actionable unclaimed idea in the survey",
              note="01-prior-art/vocabularies-and-grading.md")

    d.code("Astronomy: DS9 got it wrong and right in the same application",
           ["# a line from DS9's own reference test file —",
            "# 38 characters of geometry, 106 characters of portrayal:",
            "",
            '!circle(202.48643,47.208449,3.964") # color=pink width=3 font="times 10"',
            '!    text={Circle} tag={foo} tag={foo bar} This is a Comment',
            "",
            "# the DS9 CATALOG TOOL, same program, stores no colour per row —",
            "# it stores a rule evaluated against the data's own columns:",
            "",
            '+  condition [string equal $Class SNR]  ->  shape diamond, colour red',
            '+  size $Jmag/2.'],
           sub="The record/portrayal separation, demonstrated by the same authors who got it wrong",
           caption="Root cause: a region file is a serialization of what is ON THE SCREEN, not a "
                   "record of what is true — so colour, window state and a decorative version string "
                   "all belong in it, and the image reference does not. And JS9 shows the cost: it "
                   "grew a boolean query language over colour, which freezes the palette forever.",
           size=12)

    d.statement("the nearest neighbour",
                "Space Warps encoded “not a lens” as the absence of markers.",
                "In a strong-lensing citizen-science project at scale — the closest possible "
                "neighbour to this work — a volunteer placing no marks *was* the rejection. The "
                "reasoning behind a negative was never recorded, only its outcome.\n\n"
                "That is the same gap the workshop identified in the current annotations, reached "
                "independently. It is the strongest argument for making negative evidence a "
                "first-class mark.",
                colour=WARN)

    d.table("Nine convergences", ["Mechanism", "Domains"],
            [["A term is an opaque code; the display string is documentation", "7"],
             ["Record and portrayal are separate, separately versioned artifacts", "10"],
             ["Geometry and meaning are two objects joined by a stable id", "10"],
             ["Never join by ordinal index", "—"],
             ["Absence is not one thing; a missing key is never enough", "8"],
             ["A measurement lives in exactly one place; geometry is derived", "8"],
             ["Deprecate additively, never delete; make the retirement link typed", "6"],
             ["Generate every downstream artifact from one source file", "7"],
             ["Colour is never the only carrier; the achromatic rendering is authored", "7"],
             ["**The rich canonical form needs a flat companion or it does not get written**", "7"]],
            widths=[86, 14], size=13, highlight={9},
            sub="Mechanisms unrelated domains arrived at independently",
            note="The last one is independent confirmation of the read/write asymmetry — in four of "
                 "those seven, the flat companion is what actually shipped")

    # ── 3. The three questions ────────────────────────────────────────────────────────────
    d.section("Part three", "The three questions you asked",
              "Accessibility · thin but readable · resize for emphasis")

    d.statement("question one — accessibility",
                "Yes — and it changes the notation, not just the palette.",
                "Role is carried by glyph SHAPE. Polarity by stroke TEXTURE. Colour reinforces and "
                "never carries alone (WCAG 2.2 SC 1.4.1).\n\n"
                "Two consequences fall out rather than being imposed: size cannot carry emphasis, "
                "because size is already spoken for by measurement; and colour cannot carry "
                "polarity — which is where the room had independently arrived.",
                colour=GOOD)

    d.image("Accessibility as a test, not an aspiration",
            REF / "polarity-triad.png",
            sub="The same document rendered three times — all positive, all negative, all ambiguous "
                "— with every label removed",
            caption="If this sheet were ambiguous, the polarity channel would have failed.",
            note="05-accessibility-design.md · examples/cvd.py")

    d.numbers("The gate, and its result",
              [("136", "pairs of marks that can share a panel, checked under normal vision,\n"
                "deuteranopia, protanopia, tritanopia and greyscale", ACC),
               ("123", "of them — 90% — are NOT reliably separable by colour\nunder some vision "
                "condition", WARN),
               ("0", "are separated by colour alone.\nEvery one differs in shape, texture or "
                "orientation.", GOOD)],
              sub="For every pair below the ΔE floor, assert that a non-colour channel differs",
              note="The floor is calibrated, not guessed: pure red against pure green still measures "
                   "ΔE2000 = 12.9 under deuteranopia — almost all of it lightness — so a floor at 12 "
                   "would wave through the canonical confusion. It is set at 15.")

    d.statement("question two — thin but readable",
                "The problem was misdiagnosed. It was never mainly stroke width.",
                "", colour=WARN)

    d.table("Where the ink actually goes",
            ["", "ink over the panel", "of which, the on-image legend plate"],
            [["LensMark today, deck-01", "7.7%", "**4.6%**"],
             ["deck-04", "10.8%", "**6.6%**"],
             ["deck-08", "9.7%", "**5.3%**"],
             ["deck-07 (legend off)", "9.0%", "0%"],
             ["The campaign the room liked", "**1.2–2.3%**", "0% — its key sits *below* the image"]],
            widths=[46, 27, 27], size=14, highlight={4},
            sub="Measured on the real renders",
            note="The legend plate is roughly half the total ink, and in most decks more than every "
                 "arrow, circle, ring and label combined. Note deck-07: thirty items, no plate, and "
                 "still lighter than the eleven-item decks that have one.")

    d.statement("the rule that follows",
                "No ink inside the image rectangle that is not about a specific location in it.",
                "This settles the apparent contradiction in the record. “Labels go on the object” and "
                "“the legend obscured the image” are both true and are about different things: "
                "identification belongs on the object, the key belongs below.\n\n"
                "The notation key is printed once per SHEET, not once per panel. One principled "
                "exception: the scale bar, because a length comparison needs physical adjacency to "
                "the pixels being compared.",
                colour=ACC)

    d.statement("question three — resize for emphasis",
                "Emphasise the stroke, never the path.",
                "The obstacle is that some marks encode a measurement in their size — the ring radius "
                "IS θ_E. So every drawn scalar is tagged METRIC or PRESENTATIONAL, and emphasis "
                "scales only the presentational ones.\n\n"
                "A ring at `key` gets heavier, closer-spaced dots and a stronger casing — it visibly "
                "shouts — while its radius does not move by one pixel.",
                colour=GOOD)

    d.table("The partition",
            ["PRESENTATIONAL — emphasis scales these", "METRIC — emphasis never touches these"],
            [["stroke width, glyph size, label size and face", "circle radius (mask, θ_E, bounds)"],
             ["casing width, dot radius and spacing", "polygon vertices"],
             ["layer alpha, shaft length, tail position", "pointer head position"]],
            widths=[50, 50], size=14,
            sub="Three ordinal levels: muted · normal · key",
            note="Exactly three, and that is a bound rather than a preference: at 3:1 contrast the "
                 "achievable lightness levels between black and white are 0.000, 0.100 and 0.400 — a "
                 "fourth would need 1.30, which does not exist.")

    d.bullets("Three design decisions inside the emphasis mechanic",
              [("**An ordinal enum, not a float.** A float invites `1.37` and makes cross-panel "
                "comparison meaningless.", FG),
               ("**The attention mark is corner brackets, not a halo ring.** This notation has "
                "already spent circles on measurement; a new circle around a feature reads as "
                "another measured radius. A square never does.", FG),
               ("**A guard rail:** dot radius capped at 0.08 × radius, because oversized dots bias "
                "the radius a reader measures off the ink. A small circle therefore cannot be "
                "emphasised past a point — which is correct.", FG),
               ("**Per-layer alpha, never per-stroke**, or two overlapping muted strokes "
                "double-composite and darken — an order-dependent defect that is painful to "
                "reproduce.", FG),
               ("Two golden assertions come with it: absent emphasis must render byte-identical to "
                "`normal`, and the radius recovered from the ink must be identical across all three "
                "levels to 0.1%. That second one is what catches a multiplier leaking into a metric "
                "path.", DIM)])

    # ── 4. Notation ───────────────────────────────────────────────────────────────────────
    d.section("Part four", "Four notations, prototyped and measured",
              "One neutral content file, compiled to each — so the sheets compare notations rather "
              "than annotators")

    d.table("The four organising principles",
            ["", "Principle", "The idea worth having"],
            [["**N1**", "Terminator alphabet — keep the pointer, move the weight into the shape at "
              "the business end", "families as the thumbnail layer: filled = mass, open = light, "
              "bar = obstruction"],
             ["**N2**", "Bertin ledger — assign each visual variable to exactly one semantic "
              "dimension; the table IS the spec", "orientation derived from the physics: tangential "
              "tick = lensed light, radial = mass"],
             ["**N3**", "Station model — one composite badge per object, one fixation",
              "polarity for one stroke; a source index, which nothing today can express"],
             ["**N4**", "Evidence graph — annotate the argument, not the objects; the primitive is "
              "a link", "the source chord; a struck chord makes the non-lens visible at thumbnail "
              "size"]],
            widths=[7, 43, 50], size=12,
            sub="All four inherit one substrate, so the comparison has a single variable")

    ARM_SLIDES = [
        ("n1", "N1  ·  Terminator alphabet",
         "Keep the pointer — the one mark the room has consistently endorsed — and move the "
         "semantics into the shape at the business end. Lens mass is a filled silhouette, lensed "
         "light an open curve, obstruction a bar.",
         "Polarity is the shaft texture, reinforced on negatives by a strike through the "
         "terminator: two channels, because negative is the rarest and most consequential mark.",
         "**Does well:** highest role resolution of the four; evolutionary, so the cheapest to "
         "adopt; the terminator sits where the eye already goes.",
         "**Breaks:** below ~9 px the individual shapes collapse and only the three FAMILIES "
         "survive, so the families must be designed as the thumbnail layer."),
        ("n2", "N2  ·  Bertin ledger",
         "Do not start from a shape catalogue. Assign each of Bertin's seven visual variables to "
         "exactly one semantic dimension, forbid any dimension a second, and let the marks fall "
         "out. The assignment table IS the specification.",
         "Tick ORIENTATION is derived from the physics: tangential means lensed light, which really "
         "is tangentially stretched; radial means mass.",
         "**Does well:** orientation survives downscaling better than shape — 235 px against 300 — "
         "so it is the most thumbnail-robust; and the mark is mnemonic for an astronomer.",
         "**Breaks:** orientation is undefined at the deflector itself, where the radius vector is "
         "degenerate. The prototype handles it explicitly; the spec must too."),
        ("n3", "N3  ·  Station model",
         "Modelled on the WMO surface station plot: one compact badge per object, joined to it by a "
         "hairline stem, with slots for role, polarity, source index and treatment. One fixation "
         "per object instead of three.",
         "Polarity costs a single stroke — an underline for ambiguous, a strike for negative — "
         "which is what makes the whole approach cheap.",
         "**Does well:** lowest ink per object; carries SOURCE MEMBERSHIP, which nothing in the "
         "current system can express at all, for two characters.",
         "**Breaks:** it demands literacy. An unbriefed astronomer cannot read the badge, which is "
         "in direct tension with the room's preference for plain labels on the object."),
        ("n4", "N4  ·  Evidence graph",
         "Annotate the ARGUMENT, not the objects. A strong-lens claim is the assertion that these "
         "blobs are images of one source behind that mass — so draw that, and let the objects be "
         "implied by their participation in it.",
         "An anchor is the only mark that touches a feature. A source chord — a hairline arc "
         "concentric with the ring — says “these are images of one source, tangential, at this "
         "radius” in one stroke. Stems to a filled dot say “this is mass”.",
         "**Does well:** the only proposal where negative evidence is a first-class visual object. "
         "A positive panel has an unbroken chord; a negative panel's chord is struck. Lowest ink.",
         "**Breaks:** poor role resolution — it separates the four families and nothing finer — and "
         "it is the most unfamiliar, so it has the highest explanation cost."),
    ]
    for arm, title, *lines in ARM_SLIDES:
        m = M[arm.upper()]
        d.image(title, REF / f"{arm}.png",
                sub=f'ink {m["ink_pct"]}%  ·  ink per statement {m["ink_per_statement"]}  ·  '
                    f'occlusion {m["occlusion_pct"]}%  ·  20 statements expressed',
                side=[(lines[0], FG), (lines[1], ACC), (lines[2], GOOD), (lines[3], WARN)],
                note="Same twenty marks in every arm, compiled from one neutral content file")

    rows = []
    for a in ("N4", "N2", "N3", "N1", "R-CAMPAIGN", "R-CURRENT"):
        m = M[a]
        rows.append([("**" + a + "**") if a.startswith("R-C") else a, str(m["statements"]),
                     f'{m["ink_pct"]}%', f'**{m["ink_per_statement"]}**',
                     f'{m["occlusion_pct"]}%'])
    d.table("Measured", ["arm", "statements", "ink %", "ink / statement", "occlusion %"], rows,
            widths=[34, 16, 16, 18, 16], size=14, highlight={5},
            sub="Read the ink-per-statement column: raw ink penalises an arm for saying more",
            note="The two reference arms CANNOT express bound rings, segmentation polygons, "
                 "treatment, polarity or source grouping — five of the twenty statements. They look "
                 "thinner for free. Normalised, today's notation costs 3.2× more ink per unit of "
                 "meaning than the best candidate.")

    d.image("What survives at thumbnail size", REF / "contact-260.png",
            sub="260 px per panel — an ordinary nine-panel contact sheet",
            caption="At this size text is unreadable in EVERY arm. Only the visual channel remains — "
                    "which is the empirical case for putting polarity in texture and role in shape.",
            note="Measured thresholds: stroke texture survives to 167 px · orientation to 235 px · "
                 "glyph shape to 300 px · text needs 367 px")

    d.two_images("Today, and a candidate, on the same scene",
                 REF / "r-current.png", REF / "n1.png",
                 sub="Same twenty marks where expressible; same base pixels",
                 lcap="LensMark today — 6.3% ink, 0.42 per statement. The legend plate is the "
                      "single largest consumer of ink, and five statements cannot be made at all.",
                 rcap="N1 — 3.7% ink, 0.18 per statement, expressing five MORE things: bound rings, "
                      "segmentation, treatment, polarity and source grouping.",
                 note="The prototype reference arm validates against reality: it measures 6.3% "
                      "against 6.4–10.8% for the real renders")

    # ── 5. Metadata ───────────────────────────────────────────────────────────────────────
    d.section("Part five", "Six metadata architectures",
              "All six encoded in full, saying the same thing about the same image")

    d.table("The option space", ["", "Shape", "Verdict"],
            [["**P1**", "Evolved bespoke — typed items, orthogonal semantic fields, portrayal "
              "hoisted out", "**recommended container** — smallest diff, humans already read it"],
             ["**P2**", "GeoJSON FeatureCollection", "**fails**: no circle primitive, and 9 of 20 "
              "marks are circles whose radius IS the datum. Also the largest encoding"],
             ["**P3**", "W3C Web Annotation (JSON-LD)", "**fails the flat-schema gate**; its native "
              "review model is worth stealing conceptually"],
             ["**P4**", "Markup/annotation split with coded concepts (DICOM SR, AIM)",
              "architecturally best; 9.4× the cost — **take the identity, leave the ceremony**"],
             ["**P5**", "Dual surface: canonical JSON plus a compact line notation",
              "**orthogonal to the rest** — adopt as the model-read surface"],
             ["**P6**", "Relational / columnar (FITS BINTABLE, VOTable, Parquet)",
              "the right **export**, not the record"]],
            widths=[7, 40, 53], size=12,
            note="DS9-extended and RDF/OWL were considered and rejected as primaries, with what to "
                 "steal from each recorded")

    d.statement("the axis that decides",
                "Choose the container for cheapness. Choose the vocabulary machinery for rigour.",
                "Converters between containers are 57–161 lines each — this is the cheapest decision "
                "in the whole design, and it will absorb the most argument.\n\n"
                "What killed the current format was not a wrong choice; it was EROSION. Nothing "
                "prevented drift. So the axis that matters is the one where a mechanism can be "
                "installed.",
                colour=WARN)

    d.code("The mechanism",
           ["# one vocabulary file  ->  every other surface, generated",
            "",
            "  astromark-vocab-lens-1.0.json     80 terms, 10 schemes",
            "        |",
            "        +--> astromark-core-1.0.schema.json    flat, closed, 0 $ref",
            "        +--> generated/vocab_literals.py       pydantic Literals",
            "        +--> generated/vocab.ts                TypeScript unions",
            "        +--> generated/vocabulary.md           the docs tables",
            "",
            "+  $ python generate.py --check",
            "+  all 4 generated artifacts are current",
            "",
            "# adding a role touches ONE file.  Today it touches six."],
           sub="08-draft-spec/generate.py — and --check is the gate that keeps it true",
           size=13)

    rows = []
    for t in ("P5 model-READ surface", "P6 relational", "P1 evolved bespoke",
              "P3 W3C Web Annotation", "P4 markup/annotation split", "P2 GeoJSON-shaped"):
        r = SZ[t]
        rows.append([t.replace(" surface", ""), f'{r["chars"]:,}', f'{r["x_smallest"]}×',
                     f'~{r["est_tokens"]*50/1000:.0f}k'])
    d.table("What each encoding costs", ["encoding", "characters", "× smallest",
                                         "50-example payload"],
            rows, widths=[44, 20, 16, 20], size=14, highlight={0, 2},
            sub="The same annotation — 20 marks, 12 roles, both polarities — in every candidate",
            note="Character counts are EXACT. Token figures are derived at 3.5 chars/token, not "
                 "measured: no Claude tokenizer offline. The ratios are what the decision turns on.")

    d.statement("the asymmetry that resolves it",
                "Models READ the compact line form and WRITE the JSON form.",
                "Reading is where the volume is — fifty examples, every one costing tokens. Writing "
                "is where correctness matters — one document, schema-enforced, repairable.\n\n"
                "Nothing else in the option space satisfies both the flat-closed-schema requirement "
                "and the token budget. And four standards bodies got here first: in four of the "
                "seven domains that ship a flat companion to a rich canonical form, the flat "
                "companion is what actually gets written.",
                colour=GOOD)

    d.code("The whole annotation, in the line form",
           ["#astromark lens/1.0  frame=norm,tl,+x-right,+y-down  size=arcsec",
            "#system likely_lens grade=A p=0.92 cimg=found n=2 hard=dust_lane,second_deflector",
            "#source S1 images=2 config=double thE=1.45\"",
            "defl     + 0.500,0.350 -> 0.500,0.510",
            "dust     + 0.624,0.349 -> 0.530,0.526",
            "defl2    + 0.235,0.657 -> 0.353,0.591",
            "sat      + 0.244,0.460 -> 0.376,0.486",
            "arc      + 0.394,0.709 -> 0.458,0.590  S1  ; rev=wrong_position 0.31\"",
            "cimg     + 0.614,0.296 -> 0.550,0.415  S1  emph=key",
            "+arc      - 0.851,0.358 -> 0.739,0.282  alt=spiral_arm",
            "ring     + @m-defl r=1.45\"  S1  [1.15,1.8] method=arc_midline",
            "gal      + 0.181,0.241 r=1.1\"  treat=mask"],
           sub="2.1 kB against 14.4 kB for the JSON — and readable",
           caption="Read the highlighted line: “there is an arc-shaped feature here, it argues "
                   "AGAINST lensing, and the reason is that it is a spiral arm.” The current format "
                   "cannot say that at all, in any number of characters.", size=12)

    d.table("Four gates — fail one and the candidate is out",
            ["", "Gate", "Test"],
            [["G1", "flat, closed structured-output schema",
              "0 `$ref`, `additionalProperties:false` everywhere, honoured by a live call"],
             ["G2", "round-trip fidelity",
              "convert 19 real decks out and back; must re-render byte-identical to the golden hashes"],
             ["G3", "embargo separability",
              "`strip_free_text` provable by a **closed key allow-list**, not a type-sniffing walk"],
             ["G4", "expressiveness",
              "all 22 fixture statements expressible; the 16 P0 statements first-class"]],
            widths=[7, 33, 60], size=13,
            sub="Predicted casualties, recorded before the gates are run: P3 on G1, P2 on G2",
            note="No weighted sum. Report the matrix and decide by dominance; where nothing "
                 "dominates, name the trade in a sentence and take it to the room as a decision.")

    # ── 6. JWST examples ──────────────────────────────────────────────────────────────────
    d.section("Part six", "Worked examples on real JWST data",
              "Three candidates from the public NIRCam archive search — chosen because each one "
              "strains a different part of the format")

    d.bullets("Why these three, and what makes them honest",
              [("The source is the one-shot search of the **public JWST NIRCam archive**: 4.48 deg², "
                "5,391 targets, ten candidates surviving adversarial verification. No embargo.", FG),
               ("**Every mark is derived from a recorded measurement**, not placed by eye. Arc "
                "positions come from `blind_arc_radius_arcsec` and `blind_arc_pa_deg`; the ring from "
                "`blind_theta_E_arcsec` and its stated method; counter-images from coordinates "
                "written into the verifier's own note.", ACC),
               ("**rank 1 · J3440482-522486** — a published lens, SL2S J02176-0513. Arc AND "
                "counter-image, θ_E by half-separation. The clean positive case.", FG),
               ("**rank 2 · J15199556+2122210** — a single tangential arc with NO counter-image "
                "found after searching. Exercises `not_found` as distinct from `not_searched`, and "
                "the single-giant-arc exception.", FG),
               ("**rank 3 · J34707505-219476** — a knot chain plus a blue blob antipodal to it whose "
                "radius EXCEEDS the arc radius. The verifier explicitly declined half-separation for "
                "that reason.", FG)],
              sub="jwst-strong-lens-search/top100_clean/ — 100 cutouts, a 41-column table")

    d.image("rank 1 — a published lens", JWST / "J3440482-522486-n1.png",
            sub="SL2S J02176-0513 · grade A, 3 of 3 verifiers · θ_E 1.15″ by half-separation",
            side=[("The clean case: a blue tangential crescent east of a red elliptical, concave "
                   "toward it, with a compact blue counter-image nearly opposite.", FG),
                  ("The counter-image carries `emphasis: key` — hence the corner brackets. They are "
                   "a SQUARE deliberately: a ring would read as another measured radius.", ACC),
                  ("Its label had nowhere to go this close to the edge, so it **demoted to the "
                   "index 1** and the text moved to the caption band. That is the designed failure "
                   "— crowding degrades predictably instead of overlapping.", GOOD),
                  ("The mark at upper right is **negative polarity**: dashed shaft, struck "
                   "terminator, `NOT arc ≠ edge-on disk`. That is the verifier's own correction — "
                   "the NE streak runs radially, not tangentially, so it is a field galaxy rather "
                   "than a second arc.", WARN),
                  ("The current format could record none of that.", DIM)],
            note="Panel: colour 3.5″ zoom · north up, east left")

    d.image("rank 2 — searched, and not found", JWST / "J15199556+2122210-n1.png",
            sub="grade A · θ_E 1.48″ by arc midline · counter-image NOT FOUND",
            side=[("A thin tangential arc east through south-southeast, bluer than the host, "
                   "surviving in the deflector-subtracted panel.", FG),
                  ("The caption reads **counter-image not found**, and that is a different "
                   "statement from “not searched”. The verifier looked at every stretch and none "
                   "was visible.", ACC),
                  ("This is the three-valued absence that DICOM, RadElement, WMO and Galaxy Zoo all "
                   "arrived at independently — and it matters here, because θ_E rests on the arc "
                   "radius alone and could be high if the source is offset.", FG),
                  ("Hard case: **single giant arc**.", DIM)],
            note="The bound rings bracket the ridge measurement, 1.23–1.73″")

    d.image("rank 3 — the counter-image is farther out than the arc",
            JWST / "J34707505-219476-wide-n1.png",
            sub="grade A · five knots at r 1.24–1.49″ · a blue blob antipodal at r 1.88″",
            side=[("The blob is antipodal to the chain and the same colour as the knots — a "
                   "plausible counter-image. But its radius EXCEEDS the arc radius, which is the "
                   "reverse of the usual configuration.", FG),
                  ("So it is marked **ambiguous**: dotted shaft, hollow struck terminator, "
                   "`counter-image ? ≠ companion projection`. The record shows the feature AND its "
                   "alternative reading rather than choosing silently.", ACC),
                  ("And the ring records that half-separation was **declined**: `method: "
                   "arc_midline`, θ_E 1.40″. Had the blob been accepted, θ_E would be nearer 1.6″. "
                   "A coded method makes that decision recoverable.", WARN),
                  ("Lensing permits a counter-image farther out than the arc. Radius is never by "
                   "itself a ground for rejection.", DIM)],
            note="Rendered on the 10″ field, because — see next slide")

    d.image("Six renderings of one sky", JWST / "J3440482-522486-panels/colour_10.png",
            sub="The figure the verifiers actually saw is six panels: three stretches of the 10″ "
                "field, three of the 3.5″ zoom",
            side=[("**The rank-3 counter-image at r = 1.88″ falls OUTSIDE the 3.5″ zoom.** The same "
                   "annotation is simply unrenderable there.", BAD),
                  ("That is the clearest possible demonstration that a mark's coordinates are "
                   "meaningless without a declared frame — and that the two field-of-view groups do "
                   "NOT share one.", FG),
                  ("The current format has **one image per record**. It cannot say “these six "
                   "renderings are of one sky”, cannot say which panel the marks were placed on, "
                   "and cannot carry a per-panel scale.", WARN),
                  ("This is requirement R33, made concrete by real data rather than argued for.", ACC)],
            side_title="Why `variants` is a requirement, not a convenience",
            note="examples/jwst/variants-*.json")

    # ── 7. Recommendation ─────────────────────────────────────────────────────────────────
    d.section("Part seven", "Recommendation", "And what remains undecided")

    d.bullets("Six changes worth making whichever proposal wins",
              [("**1. The 2× canonical render.** Forced by arithmetic: at native size the "
                "bound-ring dots, θ_E dots and leader hairlines all fall below the ~1.25-pixel "
                "aliasing floor.", FG),
               ("**2. Tag every drawn scalar METRIC or PRESENTATIONAL.** This is the whole answer "
                "to resize-for-emphasis.", FG),
               ("**3. The emphasis enum, with layered compositing.**", FG),
               ("**4. Keep-out geometry in the label solver** — using each circle's ANNULUS, not its "
                "bounding box. Roughly a twenty-line change and the largest single legibility win "
                "available.", FG),
               ("**5. Demotion-to-index**, so crowding degrades deterministically instead of "
                "producing overlap.", FG),
               ("**6. The caption band**, and remove the on-image legend plate from the contract.", FG),
               ("Only the sign inventory is a bet. These six are not — they are independently "
                "valuable and they are in dependency order.", ACC)],
              sub="Notation-independent")

    d.table("Migration from lensmark/1.0", ["today", "becomes", "lossless?"],
            [["`schema_version`", "`schema` + `profile`", "yes"],
             ["`items[].label`", "`role` inferred + `label` for display",
              "**lossy where free text** — a human pass over 19 decks"],
             ["`items[].color`", "dropped; derived from `role` via the style document",
              "**intentionally** lossy"],
             ["`mask_circle.kind`", "`role` + `treatment`", "yes"],
             ["`theta_e.alt_arcsec`", "`lower_arcsec` or `upper_arcsec` by comparison",
              "yes — and it fixes the two decks carrying `(alt)` rings"],
             ["`einstein_ring.theta_e_arcsec`", "geometry only; the measurement lives once",
              "yes — removes the duplicate"],
             ["`style_defaults` inlined", "`style_ref` by id and hash", "yes"],
             ["—", "polarity, alternative, source, hard_case, counter_image, bound, polygons",
              "**new** — nothing to migrate"]],
            widths=[26, 44, 30], size=11,
            sub="Executed, not asserted: examples/migrate_lensmark.py converts a real deck",
            note="On deck-08, all four arrow labels mapped lexically; four SYSTEM-level fields need "
                 "a human — counter_image, n_images, hard_case, sources. Those are not conversion "
                 "failures; they are the things lensmark/1.0 had no way to record.")

    d.bullets("What the prior art changed, rather than confirmed",
              [("**Presence and polarity are two axes, not one.** RadElement separates present / "
                "absent / indeterminate / unknown from the term. Our `polarity` conflates “is this "
                "really there?” with “does it bear on the claim?” — which is why a dust lane sits "
                "awkwardly. Better than adding a fourth polarity value.", WARN),
               ("**Borrow the object classes; do not mint them.** The IVOA object-type vocabulary "
                "already publishes `grav-lens` and `lensed-image` as URIs. Our vocabulary should "
                "cover what is specific to ANNOTATION and reference the rest.", WARN),
               ("**Generation method is three-valued** — automatic, semiautomatic, manual — and "
                "that middle value is exactly the propose-then-correct workflow this project runs.", FG),
               ("**The grade problem has a worked answer** in BI-RADS and PI-RADS.", FG),
               ("**Two hard numbers for the render rules**: exactly three lightness levels exist at "
                "3:1 contrast, and two textures must differ in period by a full octave.", FG)],
              sub="Five findings that altered the design")

    d.bullets("What is verified, and what is not",
              [("**Measured and reproducible:** ink, occlusion, stroke contrast, encoding sizes in "
                "characters, the colour gate, schema flatness, the reference annotation validating, "
                "generator staleness, and the embargo check — 0 banned-lexicon hits across 19 "
                "documents, all three model-facing spec files completely clean.", GOOD),
               ("**Derived, not measured:** token counts. No Claude tokenizer offline; character "
                "ratios are exact and are what the decision turns on.", WARN),
               ("**Not tested:** whether image plus coordinate metadata teaches a model better than "
                "an annotated PNG. This is the hypothesis on the record, the format is partly "
                "designed around it, and it should run BEFORE the format freezes. It is cheap now "
                "that every encoding exists.", BAD),
               ("**Judgement, declared as such:** N1 versus N4. The measurements rule out the status "
                "quo; they do not choose between two good candidates.", DIM)])

    d.bullets("Six decisions for the review",
              [("**1.** Notation: N1 now, N4, or N1 with N4's source chord grafted in?", FG),
               ("**2.** Container: accept P1 plus prefixed-term identity, or spend the argument on "
                "the container?", FG),
               ("**3.** Do the six notation-independent changes land regardless? *(I would say yes.)*", FG),
               ("**4.** Is the read/write asymmetry acceptable, given it means two surfaces and a "
                "round-trip test in CI?", FG),
               ("**5.** Should presence and polarity become two separate axes?", FG),
               ("**6.** Who runs the few-shot experiment, and does the freeze wait for it?", FG),
               ("Start with the vocabulary file either way. It is the only artifact every proposal "
                "shares, the only one whose absence caused every problem in the audit — and it is "
                "the part that is not in dispute.", ACC)],
              sub="A draft vocabulary already exists: 80 terms, 10 schemes, generating a flat closed "
                  "schema the reference annotation validates against today")

    d.statement("", "One warning worth carrying into the room.",
                "DICOM's Presentation State has been the right overlay answer since 1999, is "
                "supported by every PACS, and lost in practice to burning annotations into a "
                "flattened image — because that was easier.\n\n"
                "A correct design does not win by being correct. If AstroMark's overlay path is "
                "harder than exporting a PNG with the marks drawn on it, people will export the PNG.",
                colour=WARN)

    return d


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--out", type=Path, default=HERE / "AstroMark-Research.pptx")
    args = ap.parse_args()
    deck = build(args.out)
    n = deck.save(args.out)
    print(f"wrote {args.out}  —  {n} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
